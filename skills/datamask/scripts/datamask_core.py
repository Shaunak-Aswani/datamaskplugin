"""
datamask_core.py - Shared detection logic for DataMask.

Used by both the v1 automatic script (datamask_docx.py) and the v2
interactive review app (datamask_review_app.py). Keeping this in one
place means the detection rules only need to be maintained once.

Format support: .docx, .xlsx, .pptx, .pdf. Detection (find_candidates)
is format-agnostic - it just works on extracted plain text. Each format
has its own extract_full_text_* / apply_mapping_to_* pair because the
underlying document object models are different (see notes on each
function, especially apply_mapping_to_pdf - PDFs are handled very
differently to the other three).

Optional per-format dependencies (openpyxl, python-pptx, pymupdf) are
imported lazily inside the relevant functions, so this module still
works for docx-only use even if those packages aren't installed.
"""

import hashlib
import random
import re
import string
import sys

from docx import Document

# ---------------------------------------------------------------------------
# Regex patterns for structured identifiers
# ---------------------------------------------------------------------------
EMAIL_RE = re.compile(r"[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}")
# Common dash/hyphen look-alikes seen in data pasted or exported from other
# systems (Excel, ERPs, PDFs) - a plain ASCII "-" is not the only character
# that shows up. Used both inside the number pattern and in the boundary
# check below, so a differently-encoded dash doesn't slip past either one.
_DASH_CHARS = r"\-\u2010\u2011\u2012\u2013\u2014\u2212"
# Requires a separator between the two main digit groups, and won't match
# in the middle of a longer digit run or a longer dash-chained code (the
# (?<![\w<dash>])/(?!\d) boundaries). This was tightened after real usage
# on financial documents showed the looser version flagging huge numbers
# of invoice/account/VAT/reference codes as "phone numbers" - e.g. bare
# digit blobs like "003300399701002", wrapped-cell text like "2504\n0018"
# read across a line break, and multi-segment codes like
# "PRO-DES-2504-0018" (the boundary rules out matching the tail end of a
# longer dash-chained code like that one - including when it's built with
# an en dash or other look-alike instead of a plain hyphen).
# A bare, unprefixed code shaped exactly like a phone number (e.g. just
# "2504-0018" with nothing else around it) is still an inherent,
# unresolvable-by-regex ambiguity - that's what the review step is for.
PHONE_RE = re.compile(
    r"(?<![\w" + _DASH_CHARS + r"])"
    r"(?:\+\d{1,3}[ ." + _DASH_CHARS + r"]?)?"
    r"(?:\(\d{2,4}\)[ ." + _DASH_CHARS + r"]?)?"
    r"\d{3,4}[ ." + _DASH_CHARS + r"]\d{3,4}"
    r"(?!\d)"
)

def _looks_like_year_range(text):
    """True for things like "2025-2026" or "2025\u20132026" - a fiscal-year
    or renewal-period range that is structurally identical to a phone
    number (4 digits, separator, 4 digits) and so would otherwise match
    PHONE_RE. Deliberately narrow: both halves must be plausible calendar
    years (1900-2099) and consecutive (or equal), which a real phone
    number essentially never is - this keeps it from suppressing genuine
    phone-shaped matches or unrelated dash-separated codes."""
    m = re.fullmatch(r"(\d{4})[ ." + _DASH_CHARS + r"](\d{4})", text)
    if not m:
        return False
    y1, y2 = int(m.group(1)), int(m.group(2))
    return 1900 <= y1 <= 2099 and 1900 <= y2 <= 2099 and 0 <= (y2 - y1) <= 1


# Singapore UEN (Unique Entity Number) shapes. Covers the three official
# ACRA formats, plus a looser 7-8-digit-plus-check-letter catch-all since
# not every "registration number" a document needs masked is a textbook
# UEN (e.g. GST reg numbers and older business registration numbers are
# shaped similarly but don't strictly match ACRA's spec). Word-boundaried
# so it doesn't grab the tail of a longer digit run.
UEN_RE = re.compile(
    r"\b("
    r"\d{4}\d{5}[A-Za-z]"          # local company, e.g. 201912345A
    r"|[A-Za-z]\d{2}[A-Za-z]{2}\d{4}[A-Za-z]"  # other entities, e.g. T09LL1234A
    r"|\d{7,8}[A-Za-z]"            # older business reg / generic reg no., e.g. 6738219C
    r")\b"
)

# Australian Business Number - always exactly 11 digits, conventionally
# grouped 2-3-3-3 with spaces (e.g. "34 009 200 686"). This shape is
# distinctive enough not to need an "ABN" label nearby to trigger on,
# same philosophy as UEN_RE above - though in practice it's almost always
# written right after the literal word "ABN" or "ACN" (which itself
# isn't masked, same as "UEN" isn't - only the number is).
ABN_RE = re.compile(r"\b\d{2}\s\d{3}\s\d{3}\s\d{3}\b")

# Street-type words used by the spaCy-based address matcher below. Split
# into "suffix" (the common English convention - "... Street", "...
# Avenue") and "prefix" (common in Malay/Indonesian, French, Spanish,
# German street naming - "Jalan ...", "Rue ...") styles, since a single
# suffix-only pattern would still be biased toward English-speaking
# countries' conventions. Not exhaustive - any language/convention not
# covered here (or an address with no recognizable street-type word at
# all, e.g. a bare building name) won't be auto-caught this way; that's
# what the known-entities list and your own reading in step 1 are for.
_STREET_SUFFIX_WORDS = {
    "road", "rd", "street", "st", "avenue", "ave", "drive", "dr", "lane", "ln",
    "place", "pl", "walk", "close", "crescent", "cres", "way", "boulevard",
    "blvd", "grove", "terrace", "park", "gardens", "garden", "hill", "view",
    "rise", "link", "loop", "quay", "green", "field", "circle", "court", "ct",
    "highway", "hwy", "parkway", "pkwy", "plaza", "square", "sq", "alley",
    "path", "trail", "row", "promenade", "parade", "esplanade", "strand",
    "arcade", "boardwalk", "circuit", "cove", "mews", "concourse",
}
_STREET_PREFIX_WORDS = {
    # Malay/Indonesian - confirmed real gap: "Lingkaran Syed Putra" (a
    # real Kuala Lumpur road, named after a person - common in Malaysian
    # street naming) went completely unrecognized because "lingkaran"
    # wasn't here, which meant the person-name PART of the street name
    # ("Syed Putra") fell back to being treated as an ordinary PERSON
    # entity - skipped entirely in address-only mode - rather than being
    # recognized as part of a street address the way "Jalan Ahmad" or
    # "Lorong Ahmad" already would be.
    "jalan", "lorong", "lebuh", "taman", "lingkaran", "persiaran",
    "lebuhraya", "susur", "medan", "kompleks", "wisma", "menara",
    "gang", "kampung", "kelurahan", "kecamatan",
    "rue", "avenue", "boulevard",  # French (also valid as English suffix words)
    "calle", "avenida", "paseo",  # Spanish
    "via", "viale", "piazza",  # Italian
    "strasse", "straße", "platz",  # German
}

# Legacy Singapore-only regex, kept as a fallback for when spaCy/its model
# isn't installed - see find_addresses_with_nlp for the preferred path.
_sg_street_type_alt = "|".join(sorted(_STREET_SUFFIX_WORDS | {"gardens?"}, key=len, reverse=True))
SG_ADDRESS_RE = re.compile(
    r"\b\d{1,4}[A-Za-z]?\s+[A-Z][A-Za-z'\.]*(?:\s+[A-Z][A-Za-z'\.]*)*\s+"
    r"(?:" + _sg_street_type_alt + r")\b"
    r"(?:\s*,?\s*#\d{1,3}-\d{1,5})?"
    r"(?:,?\s*Singapore(?:\s+\d{6})?)?",
    re.IGNORECASE,
)


NER_LABEL_TO_CATEGORY = {
    "PERSON": "Individual",
    "ORG": "Entity",
    "GPE": "Address",  # a bare country/city/state mention - low-risk on its own
}
# FAC/LOC (a specific named street, building, or landmark) are handled
# separately as "Location" rather than folded into NER_LABEL_TO_CATEGORY's
# "Address" - see run_spacy_pass for why: unlike a bare GPE mention (which
# scramble_address can safely pass through, since a country/city name on
# its own isn't very sensitive), a bare FAC/LOC name has no city/country
# attached for scramble_address to extract, and passing THAT through
# unchanged would leave a specific identifying place name unmasked. It
# gets a normal counted placeholder ("Location 1") instead.

# Bare corporate-suffix words (case-insensitive, trailing punctuation
# ignored) that spaCy occasionally tags as their own standalone ORG
# entity - a known NER mis-segmentation, not a real independent entity.
# This happens most often in dense tabular data (a comparables list full
# of "... CO., LTD" style names seems to confuse the boundary detector
# into sometimes splitting off just the suffix). A confirmed real bug:
# if such a fragment is accepted as its own candidate, it gets masked
# EVERYWHERE that exact text occurs - including as an unrelated fragment
# inside a company name that's already been correctly masked in full via
# a different (e.g. bare trigger-word) substitution, producing a
# nonsensical partial result like "Azure Clean Energy Co., Entity 19."
# instead of "Azure Clean Energy Co., Ltd." Filtered out entirely rather
# than trying to reattach it to a neighboring entity, since a standalone
# suffix like this carries no identifying information worth masking on
# its own anyway.
_BARE_CORPORATE_SUFFIXES = {
    "ltd", "co", "inc", "corp", "pte", "plc", "llc", "llp", "gmbh", "sa",
    "ag", "bhd", "sdn", "kk", "oy", "ab", "as", "nv", "bv", "spa", "srl",
    "limited", "incorporated", "corporation", "company", "holdings",
}

# Common building/facility-name suffix words. Used to catch a confirmed
# real spaCy mistake: a building name (often "[Company/Proper Noun]
# Tower/Centre/Plaza" - e.g. "Manulife Tower") occasionally gets tagged
# PERSON instead of an org/location label, especially when the leading
# word looks like it could be a surname or brand name to the model. Left
# as PERSON, it would get a randomized fake HUMAN name via scramble_name -
# not just wrong, but actively confusing in the masked output (a building
# reference turning into what reads as a person's name). Re-routed to
# "Location" instead whenever a PERSON-tagged span's last word is one of
# these - a real person's name essentially never ends this way.
_BUILDING_NAME_WORDS = {
    "tower", "towers", "building", "centre", "center", "plaza", "house",
    "mall", "complex", "mansion", "arcade", "point", "square", "annexe",
    "annex", "block", "wing", "podium", "hub", "quarter",
}


def _looks_like_building_name(text):
    """True if text's last word is a common building/facility-name
    suffix - see _BUILDING_NAME_WORDS docstring above for why this
    matters (a spaCy PERSON-mistag correction)."""
    words = text.strip().split()
    if not words:
        return False
    last_word = words[-1].strip(".,()[]").lower()
    return last_word in _BUILDING_NAME_WORDS


# Ordinary English words that are conventionally CAPITALIZED as defined
# terms in contracts and formal documents ("the Vessel", "the Company",
# "the Charterers", see Background above) - confirmed severe real bug:
# that capitalization convention is exactly the same cue spaCy's NER
# (and often its POS tagger too) uses to detect a genuine proper noun,
# so a term like "Vessel" in a maritime charter gets mistagged as its
# own distinct ORG/PERSON/location entity. Because every occurrence of
# identical text shares one placeholder, a SINGLE bad detection anywhere
# in the document is enough to mask EVERY legitimate occurrence of an
# extremely common word throughout the whole thing - for a word like
# "Vessel" that can be dozens of instances, not a rare one-off miss, so
# this is worth a deliberately broad, curated list rather than trying to
# special-case each word as it's discovered. Covers common contract-
# party roles, maritime/shipping terms, document-structure words, and
# generic business nouns - the categories that come up constantly across
# the kinds of legal/financial documents this tool processes.
_COMMON_DEFINED_TERMS = {
    # maritime / shipping
    "vessel", "vessels", "charter", "charterer", "charterers", "owner",
    "owners", "master", "crew", "cargo", "port", "voyage", "freight",
    "hire", "ship", "fleet", "berth", "bunkers", "flag",
    # contract parties / roles
    "company", "companies", "agreement", "contract", "party", "parties",
    "buyer", "seller", "purchaser", "vendor", "lessor", "lessee", "tenant",
    "landlord", "licensor", "licensee", "borrower", "lender", "guarantor",
    "indemnitor", "indemnitee", "assignor", "assignee", "grantor",
    "grantee", "trustee", "beneficiary", "manager", "director", "officer",
    "shareholder", "member", "client", "customer", "supplier",
    "contractor", "employer", "employee",
    # document structure
    "background", "annex", "annexure", "appendix", "schedule", "exhibit",
    "attachment", "recital", "preamble", "clause", "section", "article",
    "part", "chapter", "definitions", "notice", "certificate",
    "instrument", "deed", "summary", "overview", "introduction",
    "conclusion", "scope", "purpose", "methodology",
    # generic business
    "property", "premises", "site", "facility", "equipment", "product",
    "service", "services", "goods", "asset", "assets", "fund", "account",
    "policy", "group", "entity", "business", "corporation", "enterprise",
    "firm", "organization", "organisation", "division", "department",
    "unit", "branch", "subsidiary", "affiliate", "holding", "holdings",
    # fiscal/reporting-period abbreviations - short capitalized
    # abbreviations like these are just as prone to standalone NER
    # mistagging as any other short capitalized token (see the
    # "Vessel"/"Background" cases above) and are extremely common in
    # financial documents ("in FY 2023", "Q1 revenue", "H2 results")
    "fy", "ytd", "qtd", "mtd", "ttm", "eoy", "cy", "q1", "q2", "q3", "q4",
    "h1", "h2", "fye",
}


def _is_common_defined_term(text):
    """True if text, stripped of punctuation/whitespace and lowercased,
    is JUST a common legal/business/shipping term frequently used as a
    CAPITALIZED DEFINED TERM in contracts - see _COMMON_DEFINED_TERMS
    docstring above for why this matters and how severe it is otherwise.
    Only matches when the ENTIRE candidate is one of these words (not
    merely contains one) - "Cyan Vessel Holdings" as a genuine invented
    company name is not rejected just because "Vessel" appears in it,
    only a bare, standalone "Vessel" is."""
    stripped = text.strip().strip(".,;:()[]\"'").lower()
    return stripped in _COMMON_DEFINED_TERMS


def _looks_like_real_place_name(ent):
    """True if every substantive token in this NER-tagged span is
    POS-tagged as a proper noun (PROPN) - spaCy's OWN part-of-speech
    tagger, a SEPARATE model component from the NER tagger that
    produced the GPE/LOC/FAC label being checked here.

    Confirmed real mismatch: a common English word used as a document
    heading or section title (e.g. "Background", "Overview", "Summary")
    gets tagged as a location by the NER component surprisingly often -
    capitalization and an isolated, sentence-less heading position are
    the same cues that indicate a proper noun, so the statistical model
    latches onto them even for a plain common noun with nothing
    location-like about its actual meaning. But the SAME document's own
    POS tagger correctly identifies these as ordinary nouns/verbs (NOUN,
    VERB, ADP, ...), not proper nouns - while every genuine place name
    tested (Perth, Taiwan, Esplanade, Cross Street, Manulife Tower)
    reliably gets tagged PROPN by that same tagger. Requiring agreement
    between the two components is a strong, low-cost, generalizable
    filter - it eliminates a whole CLASS of "capitalized common word
    used as a heading" false positives, not just one specific word one
    at a time."""
    substantive = [tok for tok in ent if not tok.is_punct and not tok.is_space]
    if not substantive:
        return False
    return all(tok.pos_ == "PROPN" for tok in substantive)


def _is_bare_corporate_suffix(text):
    """True if text, once stripped of surrounding punctuation/whitespace
    and lowercased, is JUST a generic corporate-suffix word with nothing
    else - e.g. "LTD" or "Co." - as opposed to a real entity name that
    happens to end with one, e.g. "Acme Ltd" (which contains more than
    just the suffix and is fine)."""
    stripped = text.strip().strip(".,;:()[]").strip().lower()
    return stripped in _BARE_CORPORATE_SUFFIXES


def _is_implausible_entity_span(text):
    """True if a spaCy-tagged entity span is too long or too irregular to
    plausibly BE a single entity name, rather than a real one.

    Confirmed real bug: in dense tabular content (a vessel specification
    table, a rate card, anything where cells get concatenated into
    running text without normal sentence structure), spaCy occasionally
    tags an entire multi-cell chunk as one giant "ORG" - e.g. a vessel
    name, its code, a location, and a day-rate figure all glued together
    ("MMA Monarch S09 Malaysia 4,500 McDermott (65 day firm)..."). That
    is not a company name; masking the whole thing would destroy
    financial figures (day rates) and unrelated identifiers that should
    stay untouched, while ALSO not actually being "the entity" the user
    wanted masked in the first place - the real entity ("MMA") is buried
    inside it and would already be caught separately by a bare
    trigger-word match or its own correctly-scoped NER tag.

    Two independent signals, either one is disqualifying:
      - A literal newline in the matched text. A genuine single entity
        name is never going to span a line break; this is the strongest,
        lowest-false-positive-rate signal available.
      - Implausible length (a real company/person name is essentially
        never this long) OR containing 3+ separate digit-groups (a
        strong sign multiple distinct fields - a code, a location, a
        figure - got glued together rather than this being one name).
    """
    if "\n" in text:
        return True
    if len(text) > 60:
        return True
    if len(re.findall(r"\d+", text)) >= 3:
        return True
    return False


# ---------------------------------------------------------------------------
# "Scramble" placeholder generators
# ---------------------------------------------------------------------------
# Unlike the generic "Category N" numbering used for entities/clients, these
# categories are replaced with a realistic-looking but fake value instead of
# a bare counter - e.g. "Mark Rober" -> "John Doe", "6738219C" -> a
# different-but-same-shaped fake reg. number, "32 Raffles Place, Singapore"
# -> "Singapore". Generation is deterministic (seeded from the original
# text), so the same input always scrambles to the same output within and
# across a batch run - this is what keeps a name or reg. number consistent
# everywhere it appears, without needing a persistent mapping store.
_FIRST_NAMES = [
    "John", "Jane", "Julian", "Jasmine", "Jason", "Julia", "Jacob", "Joanna",
    "Justin", "Jocelyn", "James", "Joyce", "Jerome", "Jenny", "Joel", "Jia",
    "Jared", "Juliet", "Jonathan", "Janet", "Jack", "Judy",
]
_LAST_NAMES = [
    "Doe", "Des", "Davis", "Dixon", "Douglas", "Dunn", "Diaz",
    "Drake", "Duncan", "Dela Cruz", "Daniels", "Dorset", "Dawson", "Devi",
    "Damar", "Din", "Dutta", "Deng", "Duval", "Dias", "Doherty",
]


def scramble_name(original_text):
    """Deterministically map a detected person's name to a different
    fake name, e.g. "Mark Rober" -> "Julian Dixon". The mapping is
    stable for a given input string (based on a hash of it), not random
    per call.

    Deliberately constrained to always pick a first name starting with
    "J" and a last name starting with "D" (John Doe, Julian Des, Jasmine
    Davis, ...) rather than drawing from an unconstrained pool of
    realistic-looking names. This is intentional: a name that could pass
    for a real one is exactly the failure mode to avoid here - anyone
    reviewing the masked document should be able to tell at a glance
    that any "J___ D___" name is a placeholder, not second-guess whether
    it's real."""
    digest = hashlib.md5(original_text.strip().lower().encode("utf-8")).hexdigest()
    h = int(digest, 16)
    first = _FIRST_NAMES[h % len(_FIRST_NAMES)]
    last = _LAST_NAMES[(h // len(_FIRST_NAMES)) % len(_LAST_NAMES)]
    return f"{first} {last}"


def scramble_uen(original_text):
    """Deterministically scramble a registration-number-shaped string,
    preserving its exact character pattern (digit stays a digit, letter
    stays a letter matching original case, everything else - spaces,
    hyphens - is left alone) so the fake value is structurally identical
    but not traceable back to the real one, e.g. "6738219C" -> "1234567A"."""
    rnd = random.Random(original_text)
    out = []
    for ch in original_text:
        if ch.isdigit():
            out.append(str(rnd.randint(0, 9)))
        elif ch.isalpha():
            pool = string.ascii_uppercase if ch.isupper() else string.ascii_lowercase
            out.append(rnd.choice(pool))
        else:
            out.append(ch)
    return "".join(out)


def hashtag_redact(original_text):
    """Replace every letter/digit with "#", leaving spacing and
    punctuation exactly as-is - e.g. "+65 8123-4567" -> "+65 ####-####",
    "201912345A" -> "##########". Used for phone numbers and
    UEN/ABN-style registration numbers: unlike a scrambled-but-plausible
    fake value, a run of "#" can never be mistaken for real data, which
    is the point - it should be immediately obvious, not just at a
    glance but even skimming, that this was redacted."""
    return "".join("#" if ch.isalnum() else ch for ch in original_text)


# Region/state abbreviations that, when they're the last word of an
# address tail, indicate a country the source text may never spell out
# literally (an Australian address conventionally ends "<City> <STATE>
# <postcode>", never the word "Australia" itself). Checked in
# scramble_address so e.g. "Perth WA 6000" collapses to "Australia"
# rather than leaking the city+state as "Perth WA".
_REGION_CODE_TO_COUNTRY = {
    "wa": "Australia", "nsw": "Australia", "vic": "Australia", "qld": "Australia",
    "sa": "Australia", "tas": "Australia", "nt": "Australia", "act": "Australia",
}


# Every country/major territory name, recognized as-is when
# scramble_address has no comma-separated tail to extract one from (see
# the fallback in scramble_address below) - keeps e.g. bare "Australia"
# or "New Zealand" looking like a country rather than falling through to
# a generic "Address" label. Built from the full list of UN member/
# observer states plus a handful of commonly-referenced territories
# (Hong Kong, Taiwan) that come up constantly in business documents but
# aren't UN member states in their own right.
_KNOWN_COUNTRY_NAMES = {name.lower(): name for name in [
    "Afghanistan", "Albania", "Algeria", "Andorra", "Angola",
    "Antigua and Barbuda", "Argentina", "Armenia", "Australia", "Austria",
    "Azerbaijan", "Bahamas", "Bahrain", "Bangladesh", "Barbados", "Belarus",
    "Belgium", "Belize", "Benin", "Bhutan", "Bolivia",
    "Bosnia and Herzegovina", "Botswana", "Brazil", "Brunei", "Bulgaria",
    "Burkina Faso", "Burundi", "Cabo Verde", "Cambodia", "Cameroon",
    "Canada", "Central African Republic", "Chad", "Chile", "China",
    "Colombia", "Comoros", "Congo", "Costa Rica", "Croatia", "Cuba",
    "Cyprus", "Czechia", "Denmark", "Djibouti", "Dominica",
    "Dominican Republic", "Ecuador", "Egypt", "El Salvador",
    "Equatorial Guinea", "Eritrea", "Estonia", "Eswatini", "Ethiopia",
    "Fiji", "Finland", "France", "Gabon", "Gambia", "Georgia", "Germany",
    "Ghana", "Greece", "Grenada", "Guatemala", "Guinea", "Guinea-Bissau",
    "Guyana", "Haiti", "Honduras", "Hong Kong", "Hungary", "Iceland",
    "India", "Indonesia", "Iran", "Iraq", "Ireland", "Israel", "Italy",
    "Ivory Coast", "Jamaica", "Japan", "Jordan", "Kazakhstan", "Kenya",
    "Kiribati", "Kosovo", "Kuwait", "Kyrgyzstan", "Laos", "Latvia",
    "Lebanon", "Lesotho", "Liberia", "Libya", "Liechtenstein", "Lithuania",
    "Luxembourg", "Macau", "Madagascar", "Malawi", "Malaysia", "Maldives",
    "Mali", "Malta", "Marshall Islands", "Mauritania", "Mauritius",
    "Mexico", "Micronesia", "Moldova", "Monaco", "Mongolia", "Montenegro",
    "Morocco", "Mozambique", "Myanmar", "Namibia", "Nauru", "Nepal",
    "Netherlands", "New Zealand", "Nicaragua", "Niger", "Nigeria",
    "North Korea", "North Macedonia", "Norway", "Oman", "Pakistan",
    "Palau", "Palestine", "Panama", "Papua New Guinea", "Paraguay",
    "Peru", "Philippines", "Poland", "Portugal", "Qatar", "Romania",
    "Russia", "Rwanda", "Saint Kitts and Nevis", "Saint Lucia",
    "Saint Vincent and the Grenadines", "Samoa", "San Marino",
    "Sao Tome and Principe", "Saudi Arabia", "Senegal", "Serbia",
    "Seychelles", "Sierra Leone", "Singapore", "Slovakia", "Slovenia",
    "Solomon Islands", "Somalia", "South Africa", "South Korea",
    "South Sudan", "Spain", "Sri Lanka", "Sudan", "Suriname", "Sweden",
    "Switzerland", "Syria", "Taiwan", "Tajikistan", "Tanzania",
    "Thailand", "Timor-Leste", "Togo", "Tonga", "Trinidad and Tobago",
    "Tunisia", "Turkey", "Turkmenistan", "Tuvalu", "Uganda", "Ukraine",
    "United Arab Emirates", "United Kingdom", "United States", "Uruguay",
    "Uzbekistan", "Vanuatu", "Vatican City", "Venezuela", "Vietnam",
    "Yemen", "Zambia", "Zimbabwe",
]}


# Continents and broad multi-country market regions - recognized the
# same way individual countries are (see _KNOWN_COUNTRY_NAMES above),
# since financial/business documents constantly report by broad region
# rather than by individual country (a "geographic market" breakdown in
# an annual report or TP document saying "South East Asia", "Middle
# East/Africa", "Europe" rather than naming every country in each).
# Confirmed real gap otherwise: these are genuine, correctly-detected
# location references (spaCy tags them LOC, same as a real address would
# be) - the bug wasn't over-detection, it was that a region name has
# nothing to simplify DOWN to once detected (there's no bigger "country"
# to extract it from, unlike "32 Raffles Place, Singapore" -> "Singapore"),
# so it fell through to the generic "Address" placeholder exactly like a
# genuine street address would, even though a broad region name is no
# more sensitive than a bare country name is.
_KNOWN_REGION_NAMES = {name.lower(): name for name in [
    "Africa", "Antarctica", "Asia", "Europe", "North America",
    "South America", "Oceania", "Americas", "Middle East", "Far East",
    "South East Asia", "Southeast Asia", "North East Asia",
    "Northeast Asia", "East Asia", "South Asia", "Central Asia",
    "Central America", "Latin America", "Caribbean", "Scandinavia",
    "Nordic", "Nordics", "Balkans", "Sub-Saharan Africa", "North Africa",
    "West Africa", "East Africa", "Southern Africa", "Asia Pacific",
    "APAC", "EMEA", "LATAM", "ANZ", "Indochina", "Indian Subcontinent",
    "British Isles", "Iberia", "Benelux", "Gulf", "GCC", "ASEAN",
]}


def _lookup_region_or_compound(text):
    """Resolve text to itself (unchanged) if it's a recognized continent/
    broad region (see _KNOWN_REGION_NAMES), OR a "/"-joined compound of
    two or more recognized countries/regions (e.g. "Australia/New
    Zealand", "Middle East/Africa") - the same "/" shorthand
    "geographic market" tables in financial disclosures commonly use.
    Returns the ORIGINAL text (not a canonical/re-cased version) so a
    compound like "Middle East/Africa" stays exactly as written rather
    than being rebuilt from its parts, since unlike a single country
    there's no one obvious canonical form to collapse a compound to.
    Returns None if it doesn't resolve."""
    normalized = re.sub(r"\s+", " ", text).strip()
    if normalized.lower() in _KNOWN_REGION_NAMES:
        return normalized
    if "/" in normalized:
        parts = [p.strip() for p in normalized.split("/") if p.strip()]
        if len(parts) >= 2 and all(
            _lookup_country(p) or p.lower() in _KNOWN_REGION_NAMES for p in parts
        ):
            return normalized
    return None


def scramble_address(original_text):
    """Collapse a detected address down to just the country if that's
    genuinely all it amounts to (e.g. "32 Raffles Place, Singapore
    048624" -> "Singapore", since a country name on its own isn't very
    sensitive) - otherwise to the literal, non-numbered placeholder
    "Address". Deliberately not "Address 1", "Address 2", etc.: every
    detected address collapses to the exact same literal word, on
    purpose - which specific document or position an address came from
    isn't meant to be inferable from the placeholder, and a plain,
    predictable "Address" is easier to visually scan for in a reviewed
    document than a series of numbered variants would be.

    This must NEVER fall through to returning any part of the original
    text unchanged - a bare single-segment address like "54 Rose Street"
    (no comma, no country) has no country to extract, and returning it,
    or any fragment of it, as-is would silently leave real address text
    unmasked.
    """
    stripped = re.sub(r"\s+", " ", original_text).strip()
    country = _lookup_country(stripped)
    if country:
        return country
    region = _lookup_region_or_compound(stripped)
    if region:
        return region

    parts = [p.strip() for p in original_text.split(",") if p.strip()]
    if len(parts) > 1:
        # Only trust "the last comma-separated segment" as a potential
        # country when there actually WAS a comma - i.e. genuinely
        # multiple parts, not just the whole address with nothing to
        # split.
        tail = re.sub(r"\d{4,6}\s*$", "", parts[-1]).strip()
        if tail:
            country = _lookup_country(tail)
            if country:
                return country
            region = _lookup_region_or_compound(tail)
            if region:
                return region
            last_word = tail.split()[-1].lower()
            if last_word in _REGION_CODE_TO_COUNTRY:
                return _REGION_CODE_TO_COUNTRY[last_word]

    if re.search(r"\bSingapore\b", original_text, re.IGNORECASE):
        return "Singapore"

    return "Address"


# Common long-form/alternate country names that don't match
# _KNOWN_COUNTRY_NAMES's short form directly - checked after stripping a
# leading "the" (English commonly prefixes some country names with it:
# "the United Kingdom", "the Netherlands", "the Philippines").
_COUNTRY_NAME_ALIASES = {
    "united states of america": "United States",
    "russian federation": "Russia",
    "republic of korea": "South Korea",
    "democratic people's republic of korea": "North Korea",
    "people's republic of china": "China",
    "republic of china": "Taiwan",
    "uae": "United Arab Emirates",
    "uk": "United Kingdom",
    "us": "United States",
    "usa": "United States",
    "u.s.a.": "United States",
    "u.k.": "United Kingdom",
}


def _lookup_country(text):
    """Resolve text to a canonical country name if it matches - directly,
    via a common long-form alias, or after stripping a leading "the" -
    otherwise None."""
    normalized = re.sub(r"\s+", " ", text).strip().lower()
    normalized = re.sub(r"^the\s+", "", normalized)
    if normalized in _KNOWN_COUNTRY_NAMES:
        return _KNOWN_COUNTRY_NAMES[normalized]
    if normalized in _COUNTRY_NAME_ALIASES:
        return _COUNTRY_NAME_ALIASES[normalized]
    return None


def should_grey_address(category, placeholder):
    """True if a detected item should render with grey styling once
    masked, rather than looking like ordinary replacement text - used
    for the apply-time "grey box"/"grey highlight" treatment.

    Only true for an address/location that resolved to the generic
    literal "Address" placeholder - NOT one that resolved to just a
    country name (see scramble_address). A bare country mention is
    meant to stay looking like ordinary, readable text; it's the
    street-level detail collapsing to the generic "Address" that should
    visually stand out as redacted.
    """
    return category in ("Address", "Location") and placeholder == "Address"


def scramble_email(original_text):
    """Deterministically map a detected email address to a different,
    plausible-looking fake one at a generic placeholder domain, e.g.
    "jane.tan@realcompany.com" -> "j.smith42@example.com" - rather than a
    bare "Email 1" counter. Stable for a given input (same hash-based
    approach as scramble_name/scramble_uen), and deliberately never
    reuses any part of the real address (not the domain, not the local
    part) since either could still be identifying on its own."""
    digest = hashlib.md5(original_text.strip().lower().encode("utf-8")).hexdigest()
    h = int(digest, 16)
    first = _FIRST_NAMES[h % len(_FIRST_NAMES)]
    last = _LAST_NAMES[(h // len(_FIRST_NAMES)) % len(_LAST_NAMES)]
    number = h % 10000
    local = f"{first.lower().replace(' ', '')}.{last.lower()}{number}"
    return f"{local}@example.com"


# Categories that get a realistic scrambled fake value (e.g. "Mark Rober"
# -> "John Doe") instead of a generic numbered placeholder ("Individual 1").
SCRAMBLE_GENERATORS = {
    "Individual": scramble_name,
    "Name": scramble_name,
    "UEN": hashtag_redact,
    "ABN": hashtag_redact,
    "Phone": hashtag_redact,
    "Address": scramble_address,
    "Location": scramble_address,
    "Email": scramble_email,
}


def is_pre_numbered_category(category_label):
    """A category that already ends in a number (e.g. "Name 1"), or is
    explicitly marked literal with a leading "=" (e.g. "=Azure"), is
    treated as a complete, specific placeholder rather than a generic
    category to auto-number (e.g. "Client" -> "Client 1", "Client 2").
    The "=" form is what lets a known-entities CSV specify an arbitrary
    exact replacement word with no number attached at all - e.g. a
    pre-agreed codename table like "Cyan,=Azure" - not just a numbered
    category."""
    return category_label.startswith("=") or bool(re.search(r"\d\s*$", category_label))


def literal_placeholder_value(category_label):
    """Strip the "=" marker (if present) to get the actual placeholder
    text to use verbatim. For an ordinary pre-numbered category like
    "Name 1" (no "=" marker), the category itself IS the placeholder, so
    this just returns it unchanged."""
    if category_label.startswith("="):
        return category_label[1:].strip()
    return category_label


def _letter_suffix(n):
    """Convert a 1-indexed counter to a spreadsheet-style letter suffix:
    1->A, 2->B, ..., 26->Z, 27->AA, 28->AB, ... Used instead of a plain
    number for generic counted placeholders ("Entity A" rather than
    "Entity 1") - different entities still need distinct placeholders
    (collapsing them all to one literal word would make a document
    describing relationships BETWEEN entities unreadable), but a bare
    "1, 2, 3..." sequence was specifically flagged as looking too much
    like it could be real, sequential data rather than an obvious
    placeholder."""
    result = ""
    while n > 0:
        n, rem = divmod(n - 1, 26)
        result = chr(65 + rem) + result
    return result


def assign_placeholders(raw_candidates):
    """Given a flat list of candidate dicts (as returned by find_candidates,
    optionally concatenated across several files in one batch so the same
    entity gets the same placeholder everywhere), decide the placeholder
    for each distinct group and return a new list of the same candidates
    with a "placeholder" key added.

    This is factored out into core (rather than living inline in the scan
    CLI command) so that any other front end - the CLI, the Streamlit web
    app, or anything else built on top of this module - assigns
    placeholders exactly the same way instead of subtly drifting apart.

    Placeholder rules, in priority order:
      1. A known-entities-list category ending in a digit (e.g. "Name 1"),
         or explicitly marked literal with a leading "=" (e.g. "=Azure"),
         is used verbatim - see is_pre_numbered_category. A "=..." row
         also displays under a "Custom" category (rather than showing the
         raw "=Azure" marker as if it were a category name), since its
         whole point is an arbitrary one-off replacement, not a category.
      2. Individual/Name/UEN/Address get a realistic scrambled value - see
         SCRAMBLE_GENERATORS.
      3. Everything else gets a generic counted placeholder ("Client A",
         "Client B", ...), counted per category across the whole batch -
         a letter suffix rather than a number (see _letter_suffix).
    """
    counters = {}
    group_to_placeholder = {}
    result = []
    for cand in raw_candidates:
        category = cand["category"]
        group_key = cand.get("group", cand["text"])
        if group_key not in group_to_placeholder:
            if cand["source"] == "list" and is_pre_numbered_category(category):
                group_to_placeholder[group_key] = literal_placeholder_value(category)
            elif category in SCRAMBLE_GENERATORS:
                group_to_placeholder[group_key] = SCRAMBLE_GENERATORS[category](cand["text"])
            else:
                counters[category] = counters.get(category, 0) + 1
                group_to_placeholder[group_key] = f"{category} {_letter_suffix(counters[category])}"
        item = dict(cand)
        item["placeholder"] = group_to_placeholder[group_key]
        if item["category"].startswith("="):
            item["category"] = "Custom"
        result.append(item)
    return result


def load_known_entities_rows(rows):
    """
    rows: iterable of (text, category) tuples, e.g. from a parsed CSV.
    Returns them sorted longest-first so multi-word entities are matched
    before a shorter substring would be.
    """
    entities = [(t.strip(), c.strip() or "Entity") for t, c in rows if t and t.strip()]
    entities.sort(key=lambda x: len(x[0]), reverse=True)
    return entities


def try_load_ner():
    """Attempt to load spaCy for auto-detection. Returns None if unavailable."""
    try:
        import spacy

        try:
            return spacy.load("en_core_web_sm")
        except OSError:
            print(
                "[warn] spaCy installed but model 'en_core_web_sm' not found. "
                "Run: python -m spacy download en_core_web_sm — NER disabled.",
                file=sys.stderr,
            )
            return None
    except ImportError:
        return None


_address_matcher_cache = {}


def _get_address_matcher(nlp):
    """Build (once per nlp instance) a spaCy Matcher that finds generic
    street-address shapes: a leading house/block number, one or more
    proper-noun/noun tokens, and a street-type word in either suffix
    position ("32 Raffles Place") or prefix position ("Jalan Besar").
    Cached because building a Matcher is cheap but re-adding the same
    patterns every call would be silly across a multi-file scan."""
    key = id(nlp)
    if key in _address_matcher_cache:
        return _address_matcher_cache[key]

    from spacy.matcher import Matcher

    matcher = Matcher(nlp.vocab)
    matcher.add("STREET_SUFFIX", [[
        {"TEXT": {"REGEX": r"^\d+[A-Za-z]?$"}},
        {"POS": "DET", "OP": "?"},  # e.g. "12 The Esplanade" - "The" sits between the number and name
        {"POS": {"IN": ["PROPN", "NOUN"]}, "OP": "*"},  # zero, not one+ - "12 The Esplanade" has no separate name word before the type-word itself
        {"LOWER": {"IN": sorted(_STREET_SUFFIX_WORDS)}},
    ]])
    matcher.add("STREET_PREFIX", [[
        {"LOWER": {"IN": sorted(_STREET_PREFIX_WORDS)}},
        {"POS": {"IN": ["PROPN", "NOUN"]}, "OP": "+"},
    ]])
    _address_matcher_cache[key] = matcher
    return matcher


def run_spacy_pass(full_text, nlp, address_only=True):
    """Single combined pass over the text with spaCy, producing every
    NER-based candidate at once: generic street addresses (via the
    Matcher in _get_address_matcher, extended through a trailing comma
    into a following city/country entity), plus - unless address_only is
    True - names (PERSON), bare organisations (ORG), and any leftover
    place entities (GPE/LOC/FAC) that weren't already folded into a
    street match.

    address_only defaults to True: spaCy's general PERSON/ORG detection
    is a genuine source of noise on real business documents - confirmed
    directly tagging generic phrases like "Third-Party", "Vessel
    Management Services", or a bare "Group"/"Time" as if they were
    distinct organisations, each getting its own confusing numbered
    "Entity N" placeholder. Addresses (this function's other job) are a
    much narrower, more reliable detection - a real street pattern, or a
    place name, is far less likely to be confused with ordinary prose
    than an arbitrary noun phrase is. Pass address_only=False to restore
    full PERSON/ORG detection - useful if a document has names/companies
    that AREN'T already covered by a known-entities CSV and you're
    willing to review more noise in exchange for catching them
    automatically.

    This used to be two separate functions making two separate calls to
    nlp() over the same text. They're combined here for a reason beyond
    efficiency: running them separately meant a GPE/LOC/FAC entity that
    had already been absorbed into a full street address (e.g. "London"
    in "221B Baker Street, London") would ALSO come back on its own from
    the generic entity pass - producing a confusing duplicate candidate
    that, worse, wouldn't get properly scrambled on its own (a bare
    street name like "Baker Street" with no comma/city attached has
    nothing for scramble_address to extract, and previously fell through
    to being left unmasked). Tracking which entities the address matcher
    already consumed and skipping them in the leftover-entity pass avoids
    both problems.

    Returns a list of candidate dicts: {"text", "category", "source", "group"}.
    """
    import spacy

    matcher = _get_address_matcher(nlp)
    candidates = []
    chunk_size = 100000
    for offset in range(0, len(full_text), chunk_size):
        chunk = full_text[offset : offset + chunk_size]
        doc = nlp(chunk)

        spans = [doc[start:end] for _, start, end in matcher(doc)]
        spans = spacy.util.filter_spans(spans)  # drop overlaps, prefer the longest match

        ents = list(doc.ents)
        ents_by_start = {ent.start: ent for ent in ents if ent.label_ in ("GPE", "LOC", "FAC")}
        consumed_ranges = set()
        for span in spans:
            # An entity spaCy tags on the exact same tokens as our street
            # match (e.g. it separately calling "Jalan Besar" a GPE) needs
            # to be marked consumed too, not just the trailing city/
            # country entity - otherwise it gets emitted a second time,
            # standalone, in the leftover-entity loop below. Tracked by
            # (start, end) token range rather than id(ent), since spaCy
            # creates a fresh Span object (new id()) every time doc.ents
            # is accessed even for the same underlying entity.
            for ent in ents:
                if ent.start < span.end and ent.end > span.start:
                    consumed_ranges.add((ent.start, ent.end))
            i = span.end
            while i < len(doc) and doc[i].is_punct:
                i += 1
            end = span.end
            if i in ents_by_start:
                ent = ents_by_start[i]
                end = ent.end
                consumed_ranges.add((ent.start, ent.end))
                # Also greedily swallow a trailing postal/zip code number
                # right after the city/country (e.g. "Singapore 048624") -
                # otherwise it's left behind, unmasked, once the address
                # match stops at the city name.
                j = end
                if j < len(doc) and re.fullmatch(r"\d{4,6}", doc[j].text):
                    end = j + 1
            elif i < len(doc) and doc[i].is_alpha and doc[i].text[:1].isupper():
                # A capitalized word alone isn't enough evidence it's a
                # city (a building name like "Manulife Tower" would also
                # match that shape) - only commit to extending through it
                # if CONFIRMED by a recognizable marker right after: an
                # ALL-CAPS state/region code (optionally + a postcode), or
                # a bare postcode-shaped number. This is what distinguishes
                # "Perth WA 6000" (confirmed - extend) from "Manulife
                # Tower" (unconfirmed - "Tower" is neither - leave alone).
                candidate_end = i + 1
                confirmed_end = None
                if candidate_end < len(doc) and re.fullmatch(r"[A-Z]{2,3}", doc[candidate_end].text):
                    confirmed_end = candidate_end + 1
                    if confirmed_end < len(doc) and re.fullmatch(r"\d{4,6}", doc[confirmed_end].text):
                        confirmed_end += 1
                elif candidate_end < len(doc) and re.fullmatch(r"\d{4,6}", doc[candidate_end].text):
                    confirmed_end = candidate_end + 1
                if confirmed_end is not None:
                    end = confirmed_end
            full_span = doc[span.start : end]
            text = full_span.text.strip().rstrip(",").strip()
            if text:
                start_char = offset + full_span.start_char
                end_char = offset + full_span.end_char
                candidates.append({
                    "text": text, "category": "Address", "source": "ner", "group": text,
                    "start": start_char, "end": start_char + len(text),
                })

        for ent in ents:
            if (ent.start, ent.end) in consumed_ranges:
                continue  # already part of a full street-address match above
            start_char = offset + ent.start_char
            end_char = offset + ent.end_char
            if ent.label_ in ("LOC", "FAC"):
                if _is_common_defined_term(ent.text):
                    continue  # e.g. "Annex", "Building" used as a generic defined term - see docstring above
                if _is_implausible_entity_span(ent.text):
                    continue
                if not _looks_like_real_place_name(ent):
                    continue  # e.g. "Background" mistagged as a place - see docstring above
                # A specific named street/landmark with no city/country
                # attached (otherwise it would've been consumed above) -
                # gets its own "Location" category rather than "Address",
                # since scramble_address has nothing to extract a
                # country from here and passing it through unchanged
                # would leave it unmasked (see NER_LABEL_TO_CATEGORY note).
                candidates.append({
                    "text": ent.text, "category": "Location", "source": "ner", "group": ent.text,
                    "start": start_char, "end": end_char,
                })
            elif ent.label_ in NER_LABEL_TO_CATEGORY:
                if _is_common_defined_term(ent.text):
                    continue  # e.g. "Vessel", "Company", "Owners" used as a generic defined term - see docstring above
                if ent.label_ == "PERSON" and _looks_like_building_name(ent.text):
                    # e.g. "Manulife Tower" mistagged as a person - see
                    # _BUILDING_NAME_WORDS docstring. Checked BEFORE the
                    # address_only skip below (unlike other PERSON/ORG
                    # detection, this correction is itself
                    # address-relevant - it's fixing a location that got
                    # mistagged as a person, not adding general name
                    # detection - so it should still fire in address-only
                    # mode). Routes to Location rather than Individual so
                    # it doesn't get a random fake human name via
                    # scramble_name.
                    candidates.append({
                        "text": ent.text, "category": "Location", "source": "ner", "group": ent.text,
                        "start": start_char, "end": end_char,
                    })
                    continue
                if ent.label_ != "GPE" and address_only:
                    continue  # address-only mode: PERSON/ORG detection is skipped entirely - see docstring
                if ent.label_ == "ORG" and _is_bare_corporate_suffix(ent.text):
                    continue  # e.g. a bare "LTD" mis-segmented off a longer name - see note above
                if _is_implausible_entity_span(ent.text):
                    continue  # multi-line/oversized tabular garbage, not a real entity - see note above
                if ent.label_ == "GPE" and not _looks_like_real_place_name(ent):
                    continue  # e.g. "Background" mistagged as a country/city - see docstring above
                category = NER_LABEL_TO_CATEGORY[ent.label_]
                entity_end_char = end_char
                entity_text = ent.text
                if ent.label_ == "GPE":
                    # A bare country/city mention not already folded into
                    # a full street match (e.g. "Singapore 048424" sitting
                    # on its own line, with no preceding "<street>,"
                    # context to trigger the address-matcher's own
                    # extension) can still have a postcode sitting right
                    # after it - confirmed real leak otherwise, since
                    # nothing else would ever catch that trailing number.
                    j = ent.end
                    if j < len(doc) and re.fullmatch(r"\d{4,6}", doc[j].text):
                        entity_text = doc[ent.start : j + 1].text
                        entity_end_char = offset + doc[j].idx + len(doc[j].text)
                candidates.append({
                    "text": entity_text, "category": category, "source": "ner", "group": entity_text,
                    "start": start_char, "end": entity_end_char,
                })

    return candidates


def extract_full_text(doc):
    """Concatenate all paragraph text in the document (body, tables, headers/footers)."""
    parts = []
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            parts.append(paragraph.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    if paragraph.text.strip():
                        parts.append(paragraph.text)
    for section in doc.sections:
        for paragraph in section.header.paragraphs:
            if paragraph.text.strip():
                parts.append(paragraph.text)
        for paragraph in section.footer.paragraphs:
            if paragraph.text.strip():
                parts.append(paragraph.text)
    return "\n".join(parts)


def _build_flexible_entity_pattern(entity_text):
    """Build a regex for a known-entity string that tolerates two common
    real-world mismatches between a CSV list and how a name actually
    appears in a document:
      - case differences ("Acme Corp" vs "ACME CORP")
      - whitespace differences ("Bed Rock" vs "BedRock")
    Word order and the words themselves must still match exactly - this
    is not fuzzy/typo-tolerant matching, just formatting-tolerant.

    Boundaries use negative lookaround for a word character, rather than
    \\b, so that this works correctly even when the entity text itself
    starts or ends with punctuation (e.g. "EQ12, Level 10," - a
    fragment ending in a comma, common in addresses split across
    spreadsheet cells). \\b requires a transition between a word char and
    a non-word char, so a pattern ending in "," followed by whitespace in
    the real text (both non-word characters, no transition) would never
    match at all - silently, with no error, just a permanent miss. A
    lookaround only checks one side and has no such failure mode: it
    still blocks matching inside an unrelated longer word (e.g. "Bed
    Rock" inside "Bedrocking"), but doesn't care what character the
    pattern itself starts or ends with.
    """
    words = entity_text.split()
    body = r"\s*".join(re.escape(w) for w in words)
    pattern = r"(?<![A-Za-z0-9_])" + body + r"(?![A-Za-z0-9_])"
    return re.compile(pattern, re.IGNORECASE)


_word_boundary_pattern_cache = {}


def _replace_whole_word(text, original, placeholder):
    """Replace all whole-word occurrences of `original` in `text` with
    `placeholder`, using the same word-boundary logic as
    _build_flexible_entity_pattern (a negative lookaround for a word
    character, not \\b - see that function's docstring for why). Case-
    SENSITIVE (unlike the PDF path, which relies on PyMuPDF's own
    case-insensitive search_for) - matching the exact-case behavior
    this replaced. Returns (new_text, changed).

    Confirmed real, severe bug otherwise: a plain substring replacement
    (str.replace, or checking "original in text") has NO word-boundary
    awareness - a short mapping key like "us" (e.g. detected as the
    country abbreviation "US" elsewhere in the document, then looked up
    case-sensitively... or even just appearing as its own short match)
    would also match and mangle the "us" buried inside "customer",
    "focus", "industry", etc., turning it into "cPLACEHOLDERtomer".
    Every detection path in this codebase already respects word
    boundaries when finding candidates in the first place; the apply-
    time replacement needs to as well, or that care is undone at the
    last step. Cached per (original) string since the same mapping
    entry gets applied across many paragraphs/cells/runs in one file.
    """
    if original not in text:
        return text, False
    pattern = _word_boundary_pattern_cache.get(original)
    if pattern is None:
        pattern = re.compile(r"(?<![A-Za-z0-9_])" + re.escape(original) + r"(?![A-Za-z0-9_])")
        _word_boundary_pattern_cache[original] = pattern
    new_text, n = pattern.subn(lambda m: placeholder, text)
    return new_text, n > 0


def find_candidates(full_text, known_entities, nlp=None, use_regex=True, address_only=True):
    """
    Scan text and return a list of unique candidate detections:
        [{"text": ..., "category": ..., "source": "list"|"regex"|"ner"}, ...]

    Deduplicated by exact text match, first-source-wins (list > regex > ner
    priority, matching the order they're scanned below).

    CRITICAL: candidates from different sources can positionally OVERLAP
    even when their literal text differs - e.g. a known-entities-list
    bare word like "MMA" is a prefix of a company's full legal name, and
    spaCy NER separately tags that whole name (e.g. "MMA Clean Energy
    Co., Ltd.") as its own ORG entity at an overlapping position. Adding
    BOTH as independent candidates used to be a real, confirmed bug: at
    apply time each gets its own redact-and-redraw pass, and for PDFs
    especially, two overlapping replacement texts end up stamped on top
    of each other (visibly garbled/unreadable output), since erasing a
    region twice is harmless but DRAWING two different replacement
    strings at overlapping positions is not.

    The fix: every match/candidate's absolute (start, end) character span
    is tracked as it's found, in strict priority order (known-entities
    list first, then email/phone/UEN regex, then spaCy/address-regex
    last) - matching how much a human curated/trusted the source. Once a
    span is claimed by a higher-priority candidate, anything from a
    LOWER-priority source that overlaps it is dropped entirely, not just
    deduplicated by text. This also protects against the same failure
    mode happening entirely WITHIN the known-entities list itself (e.g.
    a CSV with both "Cyan" and "Cyan Renewables Pte Ltd" as separate
    rows) by trying longer entity strings first.

    use_regex: if False, skips the email/phone pattern-matching pass
    entirely (steps 2-3 below). Useful for people who want to rely purely
    on their known-entities list and avoid pattern-matching false
    positives (invoice numbers, reference codes, date ranges, etc.)
    rather than reviewing/unticking them each time.

    address_only: if True (the default), spaCy is only used for address
    detection - its general PERSON/ORG detection is skipped entirely.
    See run_spacy_pass's docstring for why this is the default: it's a
    confirmed, meaningful source of noise (generic phrases like
    "Third-Party" or "Vessel Management Services" mistagged as
    organisations) on real business documents, especially ones with
    dense tabular or legal-boilerplate content. Set to False to restore
    full name/organisation auto-detection, at the cost of more noise to
    review - worth it if a document has names/companies that AREN'T
    already covered by your known-entities CSV.
    """
    seen = {}  # text -> candidate dict
    claimed = []  # list of (start, end) spans already claimed, in priority order

    def overlaps_claimed(start, end):
        return any(start < e and end > s for s, e in claimed)

    # 1. Known entities list (highest priority). Longest entity text
    # tried first, so a more specific multi-word entity claims its span
    # before a shorter one (e.g. a bare trigger word that's also a
    # prefix of it) gets a chance to match inside it - see the docstring
    # above. Matching is case- and whitespace-flexible (see
    # _build_flexible_entity_pattern), but the candidate's "text" is
    # always the exact substring as it actually appears in the document -
    # not the CSV's spelling - since that's what has to be found again
    # later to mask it. Every variant is tagged with a shared "group" key
    # (the normalized entity string) so the caller can give "Vuchi Media
    # Inc" and "VUCHI MEDIA INC" the *same* placeholder number instead of
    # treating them as two entities.
    for entity_text, category in sorted(known_entities, key=lambda e: -len(e[0])):
        group_key = "list:" + " ".join(entity_text.split()).lower()
        pattern = _build_flexible_entity_pattern(entity_text)
        for match in pattern.finditer(full_text):
            start, end = match.start(), match.end()
            if overlaps_claimed(start, end):
                continue
            claimed.append((start, end))
            matched_text = match.group(0)
            if matched_text not in seen:
                seen[matched_text] = {
                    "text": matched_text,
                    "category": category,
                    "source": "list",
                    "group": group_key,
                }

    # 2. Emails
    if use_regex:
        for match in EMAIL_RE.finditer(full_text):
            start, end = match.start(), match.end()
            if overlaps_claimed(start, end):
                continue
            claimed.append((start, end))
            original = match.group(0)
            if original not in seen:
                seen[original] = {"text": original, "category": "Email", "source": "regex", "group": original}

    # 3. Phone numbers (7+ digits to reduce false positives on IDs/numbers)
    if use_regex:
        for match in PHONE_RE.finditer(full_text):
            original = match.group(0).strip()
            digit_count = sum(c.isdigit() for c in original)
            if digit_count < 7 or _looks_like_year_range(original):
                continue
            start, end = match.start(), match.end()
            if overlaps_claimed(start, end):
                continue
            claimed.append((start, end))
            if original not in seen:
                seen[original] = {"text": original, "category": "Phone", "source": "regex", "group": original}

    # 3b. UEN / registration numbers
    if use_regex:
        for match in UEN_RE.finditer(full_text):
            start, end = match.start(), match.end()
            if overlaps_claimed(start, end):
                continue
            claimed.append((start, end))
            original = match.group(0)
            if original not in seen:
                seen[original] = {"text": original, "category": "UEN", "source": "regex", "group": original}

    # 3b-ii. ABN (Australian Business Number) - kept as its own category
    # (still routed through the same format-preserving scramble_uen as
    # UEN, via SCRAMBLE_GENERATORS) rather than folded into the UEN regex
    # itself, since the two shapes are distinctly regional and keeping
    # them separate makes the review screen clearer about what was found.
    if use_regex:
        for match in ABN_RE.finditer(full_text):
            start, end = match.start(), match.end()
            if overlaps_claimed(start, end):
                continue
            claimed.append((start, end))
            original = match.group(0)
            if original not in seen:
                seen[original] = {"text": original, "category": "ABN", "source": "regex", "group": original}

    # 3c/4. spaCy: generic street-address detection, names, orgs, and
    # remaining locations, all in a single combined pass (see
    # run_spacy_pass docstring for why address detection and general NER
    # are combined rather than two separate scans). This is the preferred
    # address-detection path whenever an NER model is loaded - it isn't
    # tied to one country's street-naming convention, unlike the old
    # regex. Falls back to the narrower Singapore-specific regex for
    # addresses only when spaCy/its model isn't installed; name/org
    # auto-detection simply isn't available in that case (the
    # known-entities list is the way to catch those without spaCy).
    if nlp is not None:
        for cand in run_spacy_pass(full_text, nlp, address_only=address_only):
            start, end = cand.get("start"), cand.get("end")
            if start is not None:
                if overlaps_claimed(start, end):
                    continue
                claimed.append((start, end))
            if cand["text"] not in seen:
                seen[cand["text"]] = cand
    elif use_regex:
        addr_matches = sorted(SG_ADDRESS_RE.finditer(full_text), key=lambda m: (m.start(), -len(m.group(0))))
        last_end = -1
        for match in addr_matches:
            if match.start() < last_end:
                continue
            start, end = match.start(), match.end()
            last_end = match.end()
            if overlaps_claimed(start, end):
                continue
            claimed.append((start, end))
            original = match.group(0).strip().rstrip(",")
            if original and original not in seen:
                seen[original] = {"text": original, "category": "Address", "source": "regex", "group": original}

    return list(seen.values())


def build_context_snippet(full_text, target, window=45):
    """Return a short snippet of surrounding text for a candidate, for display in the review UI."""
    idx = full_text.find(target)
    if idx == -1:
        return target
    start = max(0, idx - window)
    end = min(len(full_text), idx + len(target) + window)
    prefix = ("…" if start > 0 else "") + full_text[start:idx]
    suffix = full_text[idx + len(target) : end] + ("…" if end < len(full_text) else "")
    return prefix, target, suffix


def apply_mapping_to_paragraph(paragraph, mapping, grey_texts=None, sorted_items=None):
    """
    Replace any mapped original text with its confirmed placeholder within a
    paragraph, merging runs into the first run if a change is made (same
    approach/limitation as v1 - see README).

    grey_texts: optional set of original strings (a subset of mapping's
    keys) that should get a grey highlight on the resulting run - used
    for addresses. Since all matched text in this paragraph gets merged
    into one run regardless (a pre-existing formatting limitation - see
    module docstring), if ANY of the replacements that fired in this
    paragraph were address ones, the WHOLE merged run is greyed, not
    just the address portion within it. This is an accepted
    simplification: precise per-substring highlighting would require
    rewriting the run-merging approach itself, and a paragraph
    containing an address usually doesn't have much else worth keeping
    un-greyed right next to it anyway.

    sorted_items: optional pre-sorted (longest-first) list of
    mapping.items(), so a caller applying this across many paragraphs
    (apply_mapping_to_document) can sort once up front rather than
    paying for it again on every single paragraph - a real, measurable
    cost at a few thousand mapping entries across a few thousand
    paragraphs. Falls back to sorting mapping.items() itself if not
    given, so this function still works standalone.
    """
    from docx.enum.text import WD_COLOR_INDEX

    grey_texts = grey_texts or set()
    full_text = "".join(run.text for run in paragraph.runs)
    if not full_text.strip():
        return False

    new_text = full_text
    changed = False
    any_grey = False
    # Longest original text first - see the matching note in
    # apply_mapping_to_pdf for why: a shorter key (e.g. bare "Singapore")
    # processed before a longer one that contains it as a substring (e.g.
    # "Singapore 048424") would mangle the text the longer key expects to
    # find verbatim, silently leaving whatever came after (a postcode,
    # etc.) exposed once the longer replacement becomes a no-op.
    items = sorted_items if sorted_items is not None else sorted(mapping.items(), key=lambda kv: -len(kv[0]))
    for original, placeholder in items:
        new_text, this_changed = _replace_whole_word(new_text, original, placeholder)
        if this_changed:
            changed = True
            if original in grey_texts:
                any_grey = True

    if not changed:
        return False

    if paragraph.runs:
        paragraph.runs[0].text = new_text
        if any_grey:
            paragraph.runs[0].font.highlight_color = WD_COLOR_INDEX.GRAY_25
        for run in paragraph.runs[1:]:
            run.text = ""
    return True


def apply_mapping_to_document(doc, mapping, grey_texts=None):
    """Apply a confirmed {original: placeholder} mapping across the whole document in place."""
    sorted_items = sorted(mapping.items(), key=lambda kv: -len(kv[0]))
    for paragraph in doc.paragraphs:
        apply_mapping_to_paragraph(paragraph, mapping, grey_texts, sorted_items)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    apply_mapping_to_paragraph(paragraph, mapping, grey_texts, sorted_items)
    for section in doc.sections:
        for paragraph in section.header.paragraphs:
            apply_mapping_to_paragraph(paragraph, mapping, grey_texts, sorted_items)
        for paragraph in section.footer.paragraphs:
            apply_mapping_to_paragraph(paragraph, mapping, grey_texts, sorted_items)
    return doc


# ---------------------------------------------------------------------------
# Excel (.xlsx)
# ---------------------------------------------------------------------------
def extract_full_text_xlsx(wb, with_locations=False):
    """Concatenate all literal string cell values across every sheet.

    Deliberately skips formula cells (cell.data_type == "f"). openpyxl
    returns the raw formula text as a string for those (e.g.
    "='Movements in Equity'!F14+-3.7252902984691E-09"), and scanning that
    text produces false positives - long digit runs inside floating-point
    literals or cell references get misread as phone numbers by the
    regex. Formulas are structural/computational, not narrative text, so
    they're excluded from detection entirely; only literal string cells
    (headers, note text, transaction descriptions, etc.) are scanned.

    Known limitation: worksheet tab names themselves are not scanned
    (e.g. a tab titled "Note 4.1.1 - Acme Corp" won't be caught) - see
    apply_mapping_to_xlsx's docstring for why renaming tabs safely is a
    separate, harder problem than masking cell text.

    with_locations: if True, also returns a list of (start, end, sheet_name)
    character ranges within the returned text, one per sheet that
    contributed any content. Used by the review UI to show which sheet a
    detected item actually came from - useful on workbooks with many tabs
    (a 20-sheet financial model, for instance) where "category: Phone"
    alone doesn't tell you where to look.
    """
    text_chunks = []
    locations = []
    pos = 0
    for ws in wb.worksheets:
        sheet_start = pos
        found_any = False
        for row in ws.iter_rows():
            for cell in row:
                if cell.data_type == "s" and isinstance(cell.value, str) and cell.value.strip():
                    if text_chunks:
                        text_chunks.append("\n")
                        pos += 1
                    text_chunks.append(cell.value)
                    pos += len(cell.value)
                    found_any = True
        if found_any:
            locations.append((sheet_start, pos, ws.title))
    full_text = "".join(text_chunks)
    if with_locations:
        return full_text, locations
    return full_text


def find_sheet_for_position(locations, index):
    """Given the (start, end, sheet_name) ranges from extract_full_text_xlsx
    and a character index into that same text, return which sheet it falls
    in (or None if out of range / locations wasn't provided)."""
    if index is None or index < 0:
        return None
    for start, end, sheet_name in locations:
        if start <= index < end:
            return sheet_name
    return None


def apply_mapping_to_xlsx(wb, mapping, grey_texts=None):
    """Apply a confirmed mapping across every literal string cell in every
    sheet, in place.

    grey_texts: optional set of original strings (a subset of mapping's
    keys) that should get a grey cell fill - used for addresses. Applied
    at the whole-cell level (openpyxl doesn't support per-character
    formatting within a single cell's text), so if a cell contains an
    address alongside other text, the entire cell is shaded, not just
    the address portion.

    Deliberately skips formula cells (cell.data_type == "f") for the same
    reason extract_full_text_xlsx does: a formula cell's value is not
    display text, it's executable formula syntax. Blindly substring-
    replacing inside a formula string (e.g. if a numeric literal inside a
    formula had been mis-detected as a candidate) would corrupt the
    formula and break the workbook's calculations. Only literal string
    cells are ever rewritten.

    Bigger limitation worth knowing: if a client/entity name appears only
    as a WORKSHEET TAB NAME (e.g. a tab titled "Vuchi Media") rather than
    in cell text, this function does not rename it, and there is no safe
    quick fix for that - renaming a sheet in openpyxl does not rewrite
    formulas on other sheets that reference it by name (e.g.
    ='Vuchi Media'!A1), so a naive rename would silently break every
    formula pointing at that sheet (#REF! errors). Fixing that properly
    requires parsing and rewriting every formula's sheet-name references
    across the whole workbook, which is a separate feature, not a
    one-line addition - flag it if tab names need masking too.
    """
    from openpyxl.styles import PatternFill

    grey_texts = grey_texts or set()
    grey_fill = PatternFill(start_color="B0B0B0", end_color="B0B0B0", fill_type="solid")
    changed_any = False
    # Sorted once for the whole workbook, not once per cell - with a
    # sheet of any real size (thousands of cells) and a few hundred to a
    # few thousand mapping entries, re-sorting per cell is real,
    # measurable, avoidable cost.
    sorted_items = sorted(mapping.items(), key=lambda kv: -len(kv[0]))
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if cell.data_type != "s" or not isinstance(cell.value, str) or not cell.value:
                    continue
                new_value = cell.value
                changed = False
                any_grey = False
                # Longest original text first - see the matching note in
                # apply_mapping_to_pdf for why.
                for original, placeholder in sorted_items:
                    new_value, this_changed = _replace_whole_word(new_value, original, placeholder)
                    if this_changed:
                        changed = True
                        if original in grey_texts:
                            any_grey = True
                if changed:
                    cell.value = new_value
                    if any_grey:
                        cell.fill = grey_fill
                    changed_any = True
    return changed_any


# ---------------------------------------------------------------------------
# PowerPoint (.pptx)
# ---------------------------------------------------------------------------
def _iter_pptx_text_frames(slide):
    """Yield every text frame on a slide, recursing into grouped shapes and table cells."""

    def _walk(shapes):
        for shape in shapes:
            if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                _walk(shape.shapes)
                continue
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                yield shape.text_frame
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        yield cell.text_frame

    yield from _walk(slide.shapes)
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        yield slide.notes_slide.notes_text_frame


def extract_full_text_pptx(prs):
    """Concatenate all paragraph text across every slide (shapes, tables, grouped
    shapes, and speaker notes).

    Known limitation: text baked into charts (chart titles/axis/data
    labels) and SmartArt is not covered, similar to the "images/logos"
    limitation already documented for docx.
    """
    parts = []
    for slide in prs.slides:
        for text_frame in _iter_pptx_text_frames(slide):
            for paragraph in text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs)
                if text.strip():
                    parts.append(text)
    return "\n".join(parts)


def _apply_mapping_to_pptx_paragraph(paragraph, mapping, grey_texts=None, sorted_items=None):
    """Same run-merging approach as the docx version, and the same limitation:
    formatting that varies within a single sentence gets flattened to the
    first run's formatting once a replacement happens.

    grey_texts: see apply_mapping_to_paragraph's docstring - same
    approximation (whole merged run greyed if ANY replacement in it was
    an address). PPTX doesn't have as clean a text-highlight API as docx
    via python-pptx, so this sets the run's font COLOR to grey instead
    of a background highlight - a slightly different visual (muted text
    rather than a grey box behind normal-colored text) but the same
    intent: visibly distinct from an ordinary, un-redacted run.

    sorted_items: see apply_mapping_to_paragraph's docstring - same
    reasoning (sort once for the whole presentation, not once per
    paragraph)."""
    from pptx.dml.color import RGBColor

    grey_texts = grey_texts or set()
    full_text = "".join(run.text for run in paragraph.runs)
    if not full_text.strip():
        return False
    new_text = full_text
    changed = False
    any_grey = False
    # Longest original text first - see the matching note in
    # apply_mapping_to_pdf for why.
    items = sorted_items if sorted_items is not None else sorted(mapping.items(), key=lambda kv: -len(kv[0]))
    for original, placeholder in items:
        new_text, this_changed = _replace_whole_word(new_text, original, placeholder)
        if this_changed:
            changed = True
            if original in grey_texts:
                any_grey = True
    if not changed:
        return False
    if paragraph.runs:
        paragraph.runs[0].text = new_text
        if any_grey:
            paragraph.runs[0].font.color.rgb = RGBColor(0x80, 0x80, 0x80)
        for run in paragraph.runs[1:]:
            run.text = ""
    return True


def apply_mapping_to_pptx(prs, mapping, grey_texts=None):
    """Apply a confirmed mapping across every slide in place."""
    sorted_items = sorted(mapping.items(), key=lambda kv: -len(kv[0]))
    for slide in prs.slides:
        for text_frame in _iter_pptx_text_frames(slide):
            for paragraph in text_frame.paragraphs:
                _apply_mapping_to_pptx_paragraph(paragraph, mapping, grey_texts, sorted_items)
    return prs


# ---------------------------------------------------------------------------
# Local OCR fallback for scanned/image-only PDF pages
# ---------------------------------------------------------------------------
# Uses Tesseract (a local, offline OCR engine - no cloud API, no network
# call, no credentials) rather than a cloud OCR service. This is a
# deliberate choice: cloud OCR (Azure Document Intelligence, Google
# Document AI, etc.) would need API keys and network access to a
# third-party endpoint, which may not be available or appropriate
# wherever this skill runs, and would send document content to another
# vendor. Tesseract runs entirely locally, at some cost in accuracy
# compared to a commercial cloud OCR engine - an intentional trade-off.
_tesseract_available_cache = None


def try_load_tesseract():
    """Check whether Tesseract is actually available (the pytesseract
    Python wrapper AND the underlying tesseract binary both need to be
    installed). Returns True/False rather than raising, matching the
    same graceful-degradation pattern as try_load_ner - callers should
    fall back to "no OCR" rather than fail outright.

    Cached after the first call - confirmed real, avoidable cost
    otherwise: pytesseract.get_tesseract_version() shells out to run
    "tesseract --version" as a subprocess EVERY call, and this function
    is called once per page during apply (to decide whether a given
    page needs the OCR fallback) - across a large batch of multi-page
    PDFs, that's potentially hundreds of unnecessary subprocess spawns
    for a fact (is Tesseract installed) that cannot change mid-run.
    """
    global _tesseract_available_cache
    if _tesseract_available_cache is not None:
        return _tesseract_available_cache
    try:
        import pytesseract

        pytesseract.get_tesseract_version()
        _tesseract_available_cache = True
    except Exception:
        _tesseract_available_cache = False
    return _tesseract_available_cache


_tesseract_lang_cache = None


def _tesseract_lang_string():
    """Build the language string to pass to Tesseract, combining every
    language pack actually installed (not just English) - e.g.
    "eng+dan+deu" if English, Danish, and German are all installed.
    Confirmed to matter in practice: a document in Danish (or any
    non-English language) OCR'd with English-only language data
    produces noticeably worse accuracy on that language's special
    characters (e.g. Danish æ/ø/å) - installing the relevant language
    pack (`apt-get install tesseract-ocr-dan`, etc.) and having this
    picked up automatically, with no code change needed, is the whole
    point of checking what's actually installed rather than hardcoding
    "eng". Falls back to "eng" alone if language detection itself fails
    for any reason, rather than erroring.
    """
    global _tesseract_lang_cache
    if _tesseract_lang_cache is not None:
        return _tesseract_lang_cache
    try:
        import pytesseract

        available = set(pytesseract.get_languages(config=""))
        available.discard("osd")  # orientation/script detection data, not a language
        _tesseract_lang_cache = "+".join(sorted(available)) if available else "eng"
    except Exception:
        _tesseract_lang_cache = "eng"
    return _tesseract_lang_cache


def ocr_page_lines(page, dpi=200, min_confidence=40):
    """Render a PDF page to an image and OCR it locally, returning
    recognized text grouped into LINES (not individual words):
    [{"text": line_text, "bbox": fitz.Rect}, ...], with bbox already
    converted from pixel coordinates back into PDF point coordinates so
    it can be used directly with PyMuPDF's drawing functions.

    Lines, not words, are the unit here - both for extraction and later
    for redaction. Word-level matching would require correctly
    re-joining a multi-word entity name across separate OCR word boxes
    with inconsistent spacing, which is failure-prone; whole-line
    redaction is coarser (it also blanks whatever else shares that line)
    but far safer - it can't leave a partial, still-readable fragment of
    a matched entity behind the way a slightly-off word boundary could.
    This mirrors the "fall back to whole-line/whole-paragraph redaction"
    approach already flagged as the right direction in the pre-existing
    OCR-text-layer known-limitation note below.

    Low-confidence words (Tesseract's own confidence score) are dropped
    entirely rather than guessed at, since a wrong OCR guess could either
    miss a real match or - worse - randomly happen to match something it
    shouldn't.
    """
    import fitz
    import pytesseract
    from PIL import Image

    zoom = dpi / 72.0
    pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
    img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
    lang = _tesseract_lang_string()
    data = pytesseract.image_to_data(img, lang=lang, output_type=pytesseract.Output.DICT)

    lines = {}
    for i in range(len(data["text"])):
        text = data["text"][i].strip()
        conf_raw = str(data["conf"][i])
        conf = int(conf_raw) if conf_raw.lstrip("-").isdigit() else -1
        if not text or conf < min_confidence:
            continue
        key = (data["block_num"][i], data["par_num"][i], data["line_num"][i])
        x, y, w, h = data["left"][i], data["top"][i], data["width"][i], data["height"][i]
        lines.setdefault(key, {"words": [], "boxes": []})
        lines[key]["words"].append(text)
        lines[key]["boxes"].append((x, y, x + w, y + h))

    results = []
    for info in lines.values():
        line_text = " ".join(info["words"])
        x0 = min(b[0] for b in info["boxes"]) / zoom
        y0 = min(b[1] for b in info["boxes"]) / zoom
        x1 = max(b[2] for b in info["boxes"]) / zoom
        y1 = max(b[3] for b in info["boxes"]) / zoom
        results.append({"text": line_text, "bbox": fitz.Rect(x0, y0, x1, y1)})
    return results


def _page_text_is_usable(text):
    """True if extracted PDF page text looks like real, usable text -
    False if it's empty OR looks like garbage.

    Confirmed real failure mode: a PDF with a broken or missing
    ToUnicode CMap (common in PDFs generated by "print to PDF" from a
    web page using a custom/subsetted font - e.g. a government tax
    portal) can have glyphs that render correctly on screen but extract
    as meaningless control-character byte sequences, not real Unicode
    text. A page like that LOOKS like it has native text (it's not
    empty), so the existing "does this page have any text at all"
    check used to decide whether to fall back to OCR would wrongly say
    yes - but that "text" cannot be reliably searched, matched, or
    trusted for detection at all. Left undetected, this is actively
    dangerous, not just unhelpful: known-entities/regex/NER matching
    against garbage bytes produces essentially arbitrary matches, and
    redacting based on those matches stamps grey/white boxes at
    effectively random positions across the visually-readable page,
    while the real, sensitive text underneath was never actually
    scanned for or masked at all. A page like this needs the exact same
    OCR fallback treatment as a genuinely scanned/image-only page -
    detected here by checking what fraction of the non-whitespace
    characters are control characters (a real word essentially never
    contains any).
    """
    if not text.strip():
        return False
    non_whitespace = [ch for ch in text if not ch.isspace()]
    if not non_whitespace:
        return False
    control_count = sum(1 for ch in non_whitespace if ord(ch) < 32)
    return (control_count / len(non_whitespace)) < 0.1


def extract_full_text_pdf(doc, use_ocr=True):
    """Concatenate extracted text across every page.

    Native PyMuPDF text extraction is tried first for every page (fast,
    accurate, no dependency). A page with no USABLE native text - either
    genuinely empty (a scanned/image-only page) or present but garbage
    (see _page_text_is_usable - a broken font/ToUnicode mapping) - falls
    back to local OCR, and only if use_ocr is True and Tesseract is
    actually available; otherwise that page contributes nothing, same
    as before OCR support existed, rather than erroring. A mixed
    document (some born-digital pages, some scanned/broken-text pages)
    gets the right treatment per page, not an all-or-nothing choice.
    """
    ocr_available = use_ocr and try_load_tesseract()
    parts = []
    for page in doc:
        text = page.get_text()
        if _page_text_is_usable(text):
            parts.append(text)
        elif ocr_available:
            ocr_text = "\n".join(line["text"] for line in ocr_page_lines(page))
            if ocr_text.strip():
                parts.append(ocr_text)
    return "\n".join(parts)


def _find_span_at_rect(page, rect):
    """Look up the actual font name/size PyMuPDF sees at a given rectangle,
    by matching against the page's text spans (the same data get_text('dict')
    exposes). Used so redacted placeholder text can match the size/weight of
    what it's replacing, instead of guessing from the box's pixel height.
    Returns a dict with 'font' and 'size', or None if nothing overlaps.
    """
    best = None
    best_overlap = 0
    for block in page.get_text("dict").get("blocks", []):
        for line in block.get("lines", []):
            for span in line.get("spans", []):
                sx0, sy0, sx1, sy1 = span["bbox"]
                ix0, iy0 = max(sx0, rect.x0), max(sy0, rect.y0)
                ix1, iy1 = min(sx1, rect.x1), min(sy1, rect.y1)
                if ix1 > ix0 and iy1 > iy0:
                    overlap = (ix1 - ix0) * (iy1 - iy0)
                    if overlap > best_overlap:
                        best_overlap = overlap
                        best = span
    return best


def _map_to_base14_font(fontname, flags=0):
    """Map an embedded PDF font name to the closest PyMuPDF base-14 font code.
    This is an approximation, not a real font match - PyMuPDF's redaction
    text stamping can only draw with base-14 fonts, so a document using a
    custom/branded font will still get a generic substitute. It gets the
    family (serif/mono/sans) and weight/slant right where it can detect them
    from the font name, which covers the large majority of business documents."""
    name = (fontname or "").lower()
    is_bold = "bold" in name or bool(flags & 2**4)
    is_italic = "italic" in name or "oblique" in name or bool(flags & 2**1)

    if "times" in name or "serif" in name or "georgia" in name or "garamond" in name:
        base = "tiro"  # Times family
        if is_bold and is_italic:
            return "tibi"
        if is_bold:
            return "tibo"
        if is_italic:
            return "tiit"
        return base
    if "courier" in name or "mono" in name or "consolas" in name:
        base = "cour"  # Courier family
        if is_bold and is_italic:
            return "cobi"
        if is_bold:
            return "cobo"
        if is_italic:
            return "coit"
        return base
    # default: Helvetica family (covers Arial, Calibri, most sans-serif)
    if is_bold and is_italic:
        return "hebi"
    if is_bold:
        return "hebo"
    if is_italic:
        return "heit"
    return "helv"


def _rect_matches_whole_words(rect, original, page_words):
    """True if `rect` (a match from page.search_for(original)) corresponds
    to one or more COMPLETE words, not a substring embedded inside a
    larger, unrelated word.

    Confirmed real, severe bug otherwise: PyMuPDF's search_for() does
    pure substring matching with NO word-boundary awareness at all - if
    a short string like "us" or "LTD" ends up as a mapping key (e.g.
    "US" detected as a country abbreviation elsewhere in the document,
    then looked up case-insensitively), search_for() will also match
    that exact sequence of letters wherever it happens to appear, mid-
    word or not - "us" inside "customer", "focus", "industry"; "LTD"
    inside some unrelated all-caps run. The result is silent, scattered
    corruption: two letters in the middle of an ordinary word get
    erased and replaced with a placeholder, breaking the word, while
    every other detection path in this codebase (the known-entities
    regex, the address matcher) already enforces a real word boundary
    and would never have matched there in the first place.

    page_words is the page's word-level bounding boxes from
    page.get_text("words") - (x0, y0, x1, y1, text, block_no, line_no,
    word_no) tuples - computed once per page and reused across every
    mapping entry, not recomputed per match. Every word whose box
    intersects the rect is joined back into reading order and compared
    against the original text: the match is only accepted if that
    original text appears in the combined word text with a non-word
    character (or nothing) on both sides - i.e. it lines up with real
    word boundaries in the actual words on the page, rather than
    landing partway through one.
    """
    import fitz

    target = re.sub(r"\s+", " ", original).strip().lower()
    if not target:
        return False
    overlapping = [w for w in page_words if fitz.Rect(w[:4]).intersects(rect)]
    if not overlapping:
        return False
    overlapping.sort(key=lambda w: (w[5], w[6], w[7]))  # block, line, word_no - reading order
    combined = " ".join(w[4] for w in overlapping).lower()
    idx = combined.find(target)
    if idx == -1:
        return False
    before_ok = idx == 0 or not combined[idx - 1].isalnum()
    after_idx = idx + len(target)
    after_ok = after_idx >= len(combined) or not combined[after_idx].isalnum()
    return before_ok and after_ok


def apply_mapping_to_pdf(doc, mapping, grey_texts=None):
    """Apply a confirmed mapping across every page, in place.

    grey_texts: optional set of original strings (a subset of mapping's
    keys) that should render as a solid grey redaction bar with the
    placeholder text in white, rather than the normal white-background/
    black-text look - used for addresses, so a redacted address is
    visually distinct at a glance from a normal text edit (matching how
    a blanked-out logo image already renders as a plain grey square,
    rather than looking like something a spellchecker might have
    changed).

    IMPORTANT - this works fundamentally differently from the docx/xlsx/
    pptx versions above, because PDF text isn't stored as editable
    "runs" or "cells" you can rewrite. Instead this:
      1. Finds every visual location where an original string appears
         (`page.search_for`), which returns rectangles, not text objects.
      2. Looks up the actual font size, family, and baseline PyMuPDF sees
         at that location (via `_find_span_at_rect`).
      3. Redacts (permanently removes) the text within each rectangle -
         erase only, no text stamped yet.
      4. Manually draws the placeholder text at the recorded baseline and
         font size - shrinking it only if it wouldn't otherwise fit in
         the original word's width (see below).

    Step 3/4 are deliberately split into two passes rather than using
    add_redact_annot's built-in text-stamping in one step. That built-in
    stamping fits text into the annotation's rectangle using its own line-
    height assumptions, and will silently shrink the font (observed
    shrinking a 12pt original down to 10.5pt even though the box had
    enough visual room) if its internal fit check doesn't like the
    rectangle's exact height. Drawing manually at a known baseline avoids
    that shrink entirely and reproduces the original size exactly.

    Consequences worth knowing before relying on this:
      - It's a redaction, not a text edit - the original content is
        actually removed from the page content stream, not just hidden.
      - Font size matches the original exactly, UNLESS the placeholder is
        wider than the original word (e.g. "ORA" -> "Onyx Ridge") - in
        that case it's shrunk just enough to fit the original word's
        width, to avoid bleeding into whatever text follows on the line.
        This was a real, confirmed bug (a longer placeholder rendered at
        full size, overlapping the next word) before this width check was
        added. A very large size mismatch can still end up small; there's
        a floor (5pt) below which it won't shrink further, and beyond
        that some overlap is possible - full text reflow isn't attempted.
      - Font *family* is only approximated to the closest of PyMuPDF's
        14 built-in fonts (Helvetica/Times/Courier, regular/bold/italic).
        A document using a distinctive branded font will get a generic
        substitute in that same weight/slant, not an exact match.
      - `search_for` matches visible text as PyMuPDF reconstructs it;
        text that's split oddly across PDF text objects (common in PDFs
        exported from design tools) can occasionally be missed even
        though `extract_full_text_pdf` picked it up as part of a larger
        text block. Always spot-check the output PDF.
    """
    import fitz

    grey_texts = grey_texts or set()
    replaced_count = 0
    # Sorted once, outside the page loop, rather than re-sorted on every
    # single page - sorting is cheap on its own, but there's no reason to
    # redo it once per page when the order never changes.
    sorted_mapping = sorted(mapping.items(), key=lambda kv: -len(kv[0]))
    for page in doc:
        draw_jobs = []
        # A cheap, case-insensitive plain-text substring pre-check before
        # ever calling the much more expensive page.search_for() below.
        # Confirmed real performance bug otherwise: search_for() does
        # real layout/glyph-position work per call, and with a large
        # batch (hundreds to thousands of distinct mapping entries) that
        # cost is paid for EVERY entry on EVERY page, even though the
        # overwhelming majority of entries don't even appear on most
        # given pages - measured at ~10 seconds for ~1,900 entries on a
        # single page, which multiplies straight through by page count
        # (a 50-page document: ~8 minutes; visible as the tool appearing
        # to hang on "Applying..." with no further feedback). This
        # pre-filter is not just a heuristic shortcut: every candidate's
        # text originally came from this same PDF's extracted text in
        # the first place (that's how it was detected during scan), so
        # if it's not present in this page's plain text, it CANNOT have
        # a match for search_for to find here either - skipping it is
        # always safe, never a missed match.
        # If this page's native text is garbage (see _page_text_is_usable -
        # a broken font/ToUnicode mapping, confirmed to happen with PDFs
        # exported from web pages using a custom font), skip the native
        # search_for pass entirely rather than let it search+redact based
        # on meaningless bytes - this page falls through to the OCR
        # fallback pass further below instead, same as a genuinely
        # scanned page. Computed once here and reused for the OCR trigger
        # below, rather than calling page.get_text() twice.
        page_raw_text = page.get_text()
        page_has_usable_text = _page_text_is_usable(page_raw_text)
        page_text_lower = page_raw_text.lower() if page_has_usable_text else ""
        # Word-level bounding boxes for this page, computed once and
        # reused for every mapping entry - used to validate that a
        # search_for() match corresponds to a genuine whole word (or
        # sequence of words), not a substring embedded inside a larger,
        # unrelated word. See _rect_matches_whole_words for why this
        # matters - it's a confirmed, real bug otherwise.
        page_words = page.get_text("words") if page_has_usable_text else []
        # Tracks rects already claimed by an earlier (higher-priority,
        # since mapping.items() iterates in the same priority order
        # find_candidates produced) mapping entry on THIS page. Needed
        # because PyMuPDF's search_for() is case-insensitive by default -
        # a later, lower-priority entry (e.g. a bare NER fragment like
        # "LTD") would otherwise also match, and re-redact-and-redraw,
        # a DIFFERENT-case occurrence already handled by an earlier entry
        # (e.g. "Ltd." inside a company name already masked via a bare
        # trigger-word substitution) - stamping two different placeholder
        # strings on top of each other at the same spot. This is the same
        # root problem the scan-time overlap fix addresses, but at apply
        # time and using page rects (the two occurrences can be totally
        # unrelated, non-overlapping positions AT SCAN time, and still
        # collide here because search_for's case-insensitivity widens
        # what counts as "the same text" beyond what scan-time saw).
        claimed_rects = []

        def _overlaps_claimed(rect):
            # Shrink slightly before checking - confirmed real bug
            # otherwise: two words on CONSECUTIVE lines can have bounding
            # boxes that touch or overlap by a hairline (a fraction of a
            # point) purely from font line-height/leading, even though
            # they share no actual content or position. Unshrunk, that
            # hairline touch registers as "intersects" and wrongly blocks
            # genuinely unrelated text on the next line from ever being
            # redacted. A substantial, real overlap (a longer match's rect
            # actually containing a shorter one, words sharing real
            # horizontal space) still triggers correctly after shrinking.
            shrunk = fitz.Rect(rect.x0 + 0.75, rect.y0 + 0.75, rect.x1 - 0.75, rect.y1 - 0.75)
            if shrunk.x1 <= shrunk.x0 or shrunk.y1 <= shrunk.y0:
                shrunk = rect  # too small to shrink meaningfully - use as-is
            return any(shrunk.intersects(r) for r in claimed_rects)

        # Longest original text first. Confirmed real bug otherwise: two
        # candidates can be correctly detected as separate, non-
        # overlapping matches at SCAN time (e.g. bare "Singapore" from one
        # part of the document, and "Singapore 048424" from a completely
        # different part) - but at APPLY time, search_for() searches the
        # whole page/document globally for each mapping key independently.
        # "Singapore" (the shorter key) would also match the "Singapore"
        # substring INSIDE "Singapore 048424" elsewhere on the page. If
        # that shorter entry is processed first, it claims that spot and
        # blocks the longer, more complete match from ever running there -
        # leaving the trailing postcode exposed. Doing longest-first means
        # the more complete match always gets first claim on any position
        # it covers.
        for original, placeholder in sorted_mapping:
            if original.lower() not in page_text_lower:
                continue
            # A search string containing a literal newline (i.e. the
            # matched text spans multiple visual LINES on the page, not
            # just multiple separate occurrences) makes search_for()
            # return one rect PER LINE for what is logically ONE match -
            # confirmed directly against PyMuPDF's actual behavior.
            # Drawing the full placeholder text at every one of those
            # rects would stamp it multiple times at different
            # y-positions (a real, observed bug) - only the first rect
            # for such a match gets the placeholder text; the rest are
            # still fully erased, just left blank rather than redrawn.
            is_multiline_match = "\n" in original
            for rect_index, rect in enumerate(page.search_for(original)):
                if _overlaps_claimed(rect):
                    continue
                if not is_multiline_match and not _rect_matches_whole_words(rect, original, page_words):
                    continue  # e.g. "us" matching inside "customer" - see docstring above
                span = _find_span_at_rect(page, rect)
                if span is not None:
                    fontsize = span["size"]
                    fontname = _map_to_base14_font(span["font"], span.get("flags", 0))
                    baseline_y = span["origin"][1]
                else:
                    # fallback if no matching span was found (rare)
                    fontsize = max(6, min(11, rect.height * 0.7))
                    fontname = "helv"
                    baseline_y = rect.y1 - rect.height * 0.2

                # If the placeholder is wider than the space the original
                # word occupied, shrink it to fit rather than let it run
                # into whatever follows on the line.
                available_width = rect.width
                try:
                    needed_width = fitz.get_text_length(placeholder, fontname=fontname, fontsize=fontsize)
                except Exception:
                    needed_width = None
                if needed_width and available_width > 0 and needed_width > available_width:
                    fontsize = max(5, fontsize * (available_width / needed_width))

                is_grey = original in grey_texts
                fill_color = (0.6, 0.6, 0.6) if is_grey else (1, 1, 1)
                page.add_redact_annot(rect, fill=fill_color)
                claimed_rects.append(rect)
                replaced_count += 1
                if is_multiline_match and rect_index > 0:
                    continue  # already drew the placeholder once, at the first line
                draw_jobs.append((rect, placeholder, baseline_y, fontsize, fontname, is_grey))
        page.apply_redactions()
        for rect, placeholder, baseline_y, fontsize, fontname, is_grey in draw_jobs:
            page.insert_text(
                (rect.x0, baseline_y),
                placeholder,
                fontsize=fontsize,
                fontname=fontname,
                color=(1, 1, 1) if is_grey else (0, 0, 0),
            )

        # OCR fallback pass - only runs for a page with NO native text at
        # all (a genuinely scanned/image-only page); a page that already
        # matched something above via search_for has real text and skips
        # this entirely. Redaction here is at whole-LINE granularity (see
        # ocr_page_lines) - coarser than the word-precise native-text path
        # above, but far safer for OCR'd content: it can't leave a partial,
        # still-legible fragment of a matched entity behind. The line is
        # redacted and redrawn with the substitution applied directly to
        # its recognized text, at a font size approximated from the box
        # height (there's no real font metadata to recover from a scanned
        # image the way _find_span_at_rect gets it from real text spans).
        if not page_has_usable_text and try_load_tesseract():
            for line in ocr_page_lines(page):
                line_text = line["text"]
                new_text = line_text
                new_text_lower = new_text.lower()
                changed = False
                line_is_grey = False
                # Longest original text first - see the matching note in
                # apply_mapping_to_pdf for why. Uses sorted_mapping
                # (computed once above) rather than re-sorting per line,
                # and skips the regex compile+search entirely for any
                # entry whose text isn't even present in this line's
                # plain text - same reasoning as the search_for pre-
                # filter above: building a fresh regex per entry per
                # line is real, avoidable cost at this scale.
                for original, placeholder in sorted_mapping:
                    if original.lower() not in new_text_lower:
                        continue
                    pattern = _build_flexible_entity_pattern(original)
                    if pattern.search(new_text):
                        new_text = pattern.sub(placeholder, new_text)
                        new_text_lower = new_text.lower()
                        changed = True
                        if original in grey_texts:
                            line_is_grey = True
                if not changed:
                    continue
                rect = line["bbox"]
                page.add_redact_annot(rect, fill=(0.6, 0.6, 0.6) if line_is_grey else (1, 1, 1))
                page.apply_redactions()
                fontsize = max(5, min(14, rect.height * 0.75))
                page.insert_text(
                    (rect.x0, rect.y1 - rect.height * 0.2),
                    new_text,
                    fontsize=fontsize,
                    fontname="helv",
                    color=(1, 1, 1) if line_is_grey else (0, 0, 0),
                )
                replaced_count += 1
    return replaced_count


# ---------------------------------------------------------------------------
# Images (logos / entity-identifying graphics)
# ---------------------------------------------------------------------------
# None of the text-extraction functions above can see text baked into an
# image (a letterhead logo, a signature block scan, a screenshot with a
# company name in it). That kind of detection needs actual vision, not
# regex/NER over extracted text - so instead of trying to OCR every image
# automatically, these functions just get images OUT of the document so
# Claude can look at them directly (it can view image files), decide which
# ones are client-identifying, and then have the flagged ones blanked out
# during apply. See SKILL.md step 1b for the intended workflow.
def extract_images_for_format(file_bytes, extension, out_dir):
    """Extract every embedded image in the document to out_dir.

    Returns a list of {"id", "path", "location"} dicts. "id" is the stable
    handle used later to tell redact_images_for_format / apply which images
    to blank - for OOXML formats (.docx/.xlsx/.pptx) it's the image's path
    inside the underlying zip (e.g. "word/media/image1.png"); for PDFs it's
    "<page_number>:<xref>".

    Known limitation: this pulls every embedded raster image, including
    ones that aren't logos or identifying at all (a chart screenshot, a
    decorative background). That's expected - filtering is the job of
    whoever (Claude, or the user) reviews the extracted images, not this
    function.
    """
    import io
    import os

    os.makedirs(out_dir, exist_ok=True)
    extension = extension.lower()
    images = []

    if extension in (".docx", ".pptx", ".xlsx"):
        import zipfile

        prefix = {".docx": "word/media/", ".pptx": "ppt/media/", ".xlsx": "xl/media/"}[extension]
        zf = zipfile.ZipFile(io.BytesIO(file_bytes))
        media_names = sorted(n for n in zf.namelist() if n.startswith(prefix))
        for name in media_names:
            data = zf.read(name)
            out_path = os.path.join(out_dir, os.path.basename(name))
            with open(out_path, "wb") as f:
                f.write(data)
            images.append({"id": name, "path": out_path, "location": name})
    elif extension == ".pdf":
        import fitz

        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page_index in range(len(doc)):
            page = doc[page_index]
            seen_xrefs = set()
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                if xref in seen_xrefs:
                    continue  # same image referenced twice on one page
                seen_xrefs.add(xref)
                try:
                    base_image = doc.extract_image(xref)
                except Exception:
                    continue
                ext = base_image.get("ext", "png")
                out_name = f"page{page_index + 1}_xref{xref}.{ext}"
                out_path = os.path.join(out_dir, out_name)
                with open(out_path, "wb") as f:
                    f.write(base_image["image"])
                images.append({
                    "id": f"{page_index}:{xref}",
                    "path": out_path,
                    "location": f"page {page_index + 1}",
                })
        doc.close()
    else:
        raise ValueError(f"Unsupported file type for image extraction: {extension}")

    return _dedupe_images_by_content(images)


def _dedupe_images_by_content(images):
    """Group images with IDENTICAL byte content into one entry, so a logo
    or letterhead repeated dozens of times (once per page/slide) shows up
    as ONE thumbnail with a count in the review screen, not dozens of
    near-identical rows to click through one by one.

    Each returned entry keeps its original "id" (the FIRST occurrence
    found) plus a new "all_ids" list (every occurrence sharing that exact
    content) and a "count". Checking the single displayed thumbnail is
    meant to flag ALL of "all_ids" for redaction, not just the one shown -
    see datamask_webapp.py and SKILL.md step 1b/3 for how "all_ids" gets
    expanded into individual redaction entries at flagging time. Images
    are grouped by exact byte-for-byte match (a hash of the file content)
    - a logo that's been re-saved/re-compressed slightly differently
    won't be caught as a duplicate, since it genuinely is different bytes;
    that's an acceptable gap, not a bug - visually-near-duplicate-but-not-
    identical images are rare for the "same letterhead on every page"
    case this exists for, and a fuzzy image-similarity comparison would
    be a much heavier, slower, and more error-prone thing to add for that
    marginal gain.
    """
    import hashlib

    groups = {}
    order = []
    for img in images:
        with open(img["path"], "rb") as f:
            digest = hashlib.md5(f.read()).hexdigest()
        if digest not in groups:
            groups[digest] = dict(img, all_ids=[img["id"]], count=1)
            order.append(digest)
        else:
            groups[digest]["all_ids"].append(img["id"])
            groups[digest]["count"] += 1
    return [groups[d] for d in order]


def _blank_placeholder_image_bytes(fmt="PNG"):
    """A plain solid-grey square used to overwrite a flagged logo/image.
    Deliberately featureless - the point is to remove the identifying
    content, not to visually match the original image (which would defeat
    the purpose if the original itself was the identifying thing)."""
    import io

    from PIL import Image

    img = Image.new("RGB", (300, 300), (211, 211, 211))
    buf = io.BytesIO()
    img.save(buf, format=fmt)
    return buf.getvalue()


def _prepare_replacement_image_bytes(replacement_image_bytes, fmt="PNG"):
    """Re-encode an arbitrary user-supplied replacement image (a company's
    real logo file, in whatever format they happened to save it in - PNG,
    JPEG, even something less common like WEBP or BMP) into the exact
    format a given embedding slot needs (PNG or JPEG, matching the
    original image's own file extension - see the fmt logic in
    redact_images_in_saved_bytes/redact_images_in_pdf). This is what makes
    "swap Client X's logo for our replacement logo everywhere" safe to do
    generically: the caller doesn't need to know or care what format the
    replacement file was saved in, or what format each individual slot in
    the document expects - this reconciles the two. JPEG has no alpha
    channel, so a replacement logo with transparency gets flattened onto
    a white background rather than failing outright, since PIL can't save
    RGBA to JPEG directly."""
    import io

    from PIL import Image

    img = Image.open(io.BytesIO(replacement_image_bytes))
    if fmt.upper() == "JPEG" and img.mode in ("RGBA", "LA", "P"):
        background = Image.new("RGB", img.size, (255, 255, 255))
        img = img.convert("RGBA")
        background.paste(img, mask=img.split()[-1])
        img = background
    elif img.mode not in ("RGB", "RGBA"):
        img = img.convert("RGBA" if "transparency" in img.info else "RGB")
    buf = io.BytesIO()
    img.save(buf, format=fmt)
    return buf.getvalue()


def redact_images_in_saved_bytes(file_bytes, extension, image_ids, replacement_image_bytes=None):
    """Post-process an already-saved .docx/.xlsx/.pptx (i.e. the bytes that
    save_document() produced) to blank out - or, if replacement_image_bytes
    is given, replace with a specific logo/image - specific embedded
    images by their zip path (the "id" extract_images_for_format
    returned).

    This runs on the SAVED bytes rather than the live python-docx/pptx/
    openpyxl object, because none of those libraries expose a clean way to
    overwrite an image part's raw bytes in place - patching the zip
    directly, once everything else has already been serialized, is far
    more reliable than fighting each library's object model for this.

    The replacement (whether the default grey square or a supplied
    replacement image) is sized to fill whatever frame the ORIGINAL image
    occupied in the document - that frame's size/position is untouched,
    only the picture data inside it changes - so if the original had a
    very different aspect ratio than the replacement, the replacement
    will appear stretched/cropped inside that original frame. This
    matters more for a real replacement logo than it did for a plain grey
    square: if a client's logo is landscape and the prepared replacement
    logo is closer to square, it may look visibly distorted - worth a
    quick visual check on the actual output, same as any other automated
    redaction. Returns (new_bytes, count_redacted).
    """
    import io
    import zipfile

    if not image_ids:
        return file_bytes, 0

    image_ids = set(image_ids)
    src = zipfile.ZipFile(io.BytesIO(file_bytes))
    out_buf = io.BytesIO()
    count = 0
    with zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED) as out:
        for item in src.infolist():
            data = src.read(item.filename)
            if item.filename in image_ids:
                ext = item.filename.rsplit(".", 1)[-1].lower()
                fmt = "JPEG" if ext in ("jpg", "jpeg") else "PNG"
                if replacement_image_bytes is not None:
                    data = _prepare_replacement_image_bytes(replacement_image_bytes, fmt)
                else:
                    data = _blank_placeholder_image_bytes(fmt)
                count += 1
            out.writestr(item, data)
    out_buf.seek(0)
    return out_buf.read(), count


def redact_images_in_pdf(doc, image_ids, replacement_image_bytes=None):
    """Blank out - or, if replacement_image_bytes is given, replace with a
    specific logo/image - specific images in an already-open PyMuPDF
    document, in place, at every location that image appears (an image
    can be referenced more than once on a page, or reused across pages).
    image_ids uses the "<page_number>:<xref>" scheme
    extract_images_for_format produces for PDFs.

    Unlike the OOXML path this doesn't need a post-save patch - PyMuPDF
    lets you draw directly onto an open page, and get_image_rects() gives
    the exact placement rectangle(s) for a given image xref on a page.
    A replacement image is inserted at that same rectangle (so it fills
    the original image's footprint, same stretch/crop caveat as the
    OOXML path above); with no replacement given, the original behavior
    (a plain opaque grey rectangle) is unchanged. Returns count of images
    redacted."""
    if not image_ids:
        return 0
    image_ids = set(image_ids)
    prepared_replacement = (
        _prepare_replacement_image_bytes(replacement_image_bytes, "PNG")
        if replacement_image_bytes is not None
        else None
    )
    count = 0
    for page in doc:
        for img in page.get_images(full=True):
            xref = img[0]
            key = f"{page.number}:{xref}"
            if key not in image_ids:
                continue
            for rect in page.get_image_rects(xref):
                if prepared_replacement is not None:
                    page.insert_image(rect, stream=prepared_replacement)
                else:
                    page.draw_rect(rect, color=(0.83, 0.83, 0.83), fill=(0.83, 0.83, 0.83))
            count += 1
    return count


# ---------------------------------------------------------------------------
# Format dispatch
# ---------------------------------------------------------------------------
def open_document(file_bytes, extension):
    """Open a document from raw bytes based on its extension (.docx/.xlsx/.pptx/.pdf)."""
    import io

    extension = extension.lower()
    if extension == ".docx":
        return Document(io.BytesIO(file_bytes))
    elif extension == ".xlsx":
        import openpyxl

        return openpyxl.load_workbook(io.BytesIO(file_bytes))
    elif extension == ".pptx":
        from pptx import Presentation

        return Presentation(io.BytesIO(file_bytes))
    elif extension == ".pdf":
        import fitz

        return fitz.open(stream=file_bytes, filetype="pdf")
    else:
        raise ValueError(f"Unsupported file type: {extension}")


def extract_text_for_format(doc, extension):
    """Dispatch to the right extractor based on file extension."""
    extension = extension.lower()
    if extension == ".docx":
        return extract_full_text(doc)
    elif extension == ".xlsx":
        return extract_full_text_xlsx(doc)
    elif extension == ".pptx":
        return extract_full_text_pptx(doc)
    elif extension == ".pdf":
        return extract_full_text_pdf(doc)
    else:
        raise ValueError(f"Unsupported file type: {extension}")


def apply_mapping_for_format(doc, mapping, extension, grey_texts=None):
    """Dispatch to the right masking function based on file extension.

    grey_texts: optional set of original strings (a subset of mapping's
    keys) to render with grey styling instead of the normal look - see
    each format-specific apply function for exactly what that means in
    that format (a solid grey box in PDF, a grey highlight in docx, grey
    text in pptx, a grey cell fill in xlsx). Used for addresses.
    """
    extension = extension.lower()
    if extension == ".docx":
        return apply_mapping_to_document(doc, mapping, grey_texts)
    elif extension == ".xlsx":
        return apply_mapping_to_xlsx(doc, mapping, grey_texts)
    elif extension == ".pptx":
        return apply_mapping_to_pptx(doc, mapping, grey_texts)
    elif extension == ".pdf":
        return apply_mapping_to_pdf(doc, mapping, grey_texts)
    else:
        raise ValueError(f"Unsupported file type: {extension}")


def save_document(doc, extension):
    """Save a document (of whatever type) to an in-memory buffer and return it."""
    import io

    buffer = io.BytesIO()
    extension = extension.lower()
    if extension == ".pdf":
        buffer.write(doc.tobytes())
    else:
        doc.save(buffer)
    buffer.seek(0)
    return buffer


def sanitize_filename_with_mapping(filename, mapping):
    """Apply the same {original: placeholder} mapping used on a document's
    content to its filename too - a client name in the document body isn't
    the only place it can leak; the original upload filename often carries
    it as well (e.g. "OnyxRidge_2026_Financials.pdf"), and that's easy to
    forget about since it never shows up in the content review step.

    Tolerates spaces, underscores, and hyphens as equivalent word
    separators (filenames commonly use underscores where document text
    uses spaces), and matches case-insensitively - same spirit as the
    document-content matching in _build_flexible_entity_pattern, but
    filenames need underscore/hyphen tolerance too, which document text
    doesn't. The placeholder's own spaces are converted to underscores in
    the result, to keep the output filename clean.

    This is still a best-effort pass, not a guarantee: it only catches
    the exact original strings already in the mapping (i.e. whatever was
    actually found and approved during content review) - it can't catch
    an abbreviation or partial name in the filename that never appeared
    verbatim in the document content itself.
    """
    result = filename
    # Longest original text first - see the matching note in
    # apply_mapping_to_pdf for why.
    for original, placeholder in sorted(mapping.items(), key=lambda kv: -len(kv[0])):
        words = original.split()
        if not words:
            continue
        # Note: deliberately NOT using \b here - regex treats underscore as
        # a "word" character, so \b fails to find a boundary right next to
        # one (e.g. "ORA_Financials" - \bORA\b would not match, since there's
        # no transition between "A" and "_"). Explicit lookarounds against
        # alphanumerics only are used instead, so underscores/hyphens count
        # as proper boundaries.
        pattern = re.compile(
            r"(?<![A-Za-z0-9])" + r"[\s_-]*".join(re.escape(w) for w in words) + r"(?![A-Za-z0-9])",
            re.IGNORECASE,
        )
        safe_placeholder = re.sub(r"\s+", "_", placeholder)
        result = pattern.sub(safe_placeholder, result)
    return result
